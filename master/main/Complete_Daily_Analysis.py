import pandas as pd
from scipy.spatial import KDTree
import numpy as np
import os

main_folder_path = r'C:\Users\Kamalesh.kb\Desktop\Lectrix_Data_Analysis_Version_1\Automationdashboard\master\main\menu_1_Daily_Analysis'

def Merge_log_km(main_folder_path):
    print("Merging Log and Km")

    # Get all project folders within the main directory
    project_folders = [os.path.join(main_folder_path, f) for f in os.listdir(main_folder_path) if os.path.isdir(os.path.join(main_folder_path, f))]

    # Process each project folder
    for project_folder in project_folders:
        all_merged_dfs = []  # List to store each merged dataframe for the current project

        # Iterate over each subdirectory in the current project folder
        for subfolder in os.listdir(project_folder):
            subfolder_path = os.path.join(project_folder, subfolder)

            # Check if the path is indeed a directory
            if os.path.isdir(subfolder_path):
                log_km_file = os.path.join(subfolder_path, 'log_km.csv')
                log_file = os.path.join(subfolder_path, 'log.csv')
                km_file = os.path.join(subfolder_path, 'km.csv')

                # Check if 'log_km.csv' does not exist and both 'log.csv' and 'km.csv' exist
                if not os.path.exists(log_km_file) and os.path.exists(log_file) and os.path.exists(km_file):
                    # Load the CSV files
                    log_df = pd.read_csv(log_file)
                    km_df = pd.read_csv(km_file)

                    # Process the files as per your existing logic
                    log_df['localtime'] = pd.to_datetime(log_df['localtime'], dayfirst=True)
                    km_df['localtime'] = pd.to_datetime(km_df['localtime'], dayfirst=True)
                    log_timestamps = log_df['localtime'].astype('int64').to_numpy()
                    km_timestamps = km_df['localtime'].astype('int64').to_numpy()

                    km_tree = KDTree(km_timestamps.reshape(-1, 1))
                    distances, indices = km_tree.query(log_timestamps.reshape(-1, 1))

                    log_df['km2_index'] = indices
                    merged_df = pd.merge(log_df, km_df, left_on='km2_index', right_index=True, suffixes=('_log', '_km'))

                    merged_df.rename(columns={
                        'id_log': 'id',
                        'timestamp_log': 'timestamp',
                        'localtime_log': 'localtime'
                    }, inplace=True)

                    merged_df.drop(columns=['km2_index'], inplace=True)

                    # Save the merged_df to the same subfolder as "log_km.csv"
                    merged_df_output_path = log_km_file
                    merged_df.to_csv(merged_df_output_path, index=False)
                    print(f"Merged file saved to {merged_df_output_path}")
                else:
                    if os.path.exists(log_km_file):
                        print(f"'log_km.csv' already exists in folder: {subfolder}")
                    if not os.path.exists(log_file):
                        print(f"'log.csv' not found in folder: {subfolder}")
                    if not os.path.exists(km_file):
                        print(f"'km.csv' not found in folder: {subfolder}")



def crop_data(main_folder_path):
    print("Crop out the anomaly")
    import os
    import pandas as pd
    import matplotlib.pyplot as plt
    import matplotlib.dates as mdates
    import mplcursors  # Import mplcursors
    import warnings

    # Disable the SettingWithCopyWarning
    warnings.filterwarnings('ignore', category=pd.errors.SettingWithCopyWarning)

    def plot_ghps(data,folder_name):
        if 'localtime' not in data.columns:
            # Convert the 'timestamp' column from Unix milliseconds to datetime
            data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')

            # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
            data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)


        # Convert 'localtime' column to datetime
        data['localtime'] = pd.to_datetime(data['localtime'], format='%Y-%m-%d %H:%M:%S.%f')
        data['localtime'] = pd.to_datetime(data['localtime'])
        data.set_index('localtime', inplace=True)
        # data['PackCurr_6'] = data.apply(adjust_current, axis=1)

        # Create a figure and axes for plotting
        fig, ax1 = plt.subplots(figsize=(10, 6))

        # Plot 'PackCurr_6' on primary y-axis
        line1, = ax1.plot(data.index, -data['PackCurr_6'], color='blue', label='PackCurr_6')
        ax1.set_ylabel('Pack Current (A)', color='blue')
        ax1.yaxis.set_label_coords(-0.1, 0.7)  # Adjust label position

        # Create secondary y-axis for 'MotorSpeed_340920578' (RPM)
        ax2 = ax1.twinx()
        line2, = ax2.plot(data.index, data['MotorSpeed_340920578'], color='green', label='Motor Speed')
        ax2.set_ylabel('Motor Speed (RPM)', color='green')

        # Add 'AC_Current_340920579' to primary y-axis
        line3, = ax1.plot(data.index, data['AC_Current_340920579'], color='lightgray', label='AC Current')

        # Add 'AC_Voltage_340920580' scaled to 10x to the left side y-axis
        line4, = ax1.plot(data.index, data['AC_Voltage_340920580'] * 10, color='lightgray', label='AC Voltage (x10)')

        # Add 'Throttle_408094978' to the left side y-axis
        line5, = ax1.plot(data.index, data['Throttle_408094978'], color='lightgray', label='Throttle (%)')

            # Add 'Throttle_408094978' to the left side y-axis
        line6, = ax1.plot(data.index, data['SOC_8'], color='red', label='SOC (%)')

        # Hide the y-axis label for 'AC_Current_340920579'
        ax1.get_yaxis().get_label().set_visible(False)

        # Set x-axis label and legend
        ax1.set_xlabel('Local Time')
        ax1.legend(loc='upper left')
        ax2.legend(loc='upper right')

        # Add a title to the plot
        plt.title('Battery Pack, Motor Data, and Throttle')

        # Format x-axis ticks as hours:minutes:seconds
        ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M:%S'))

        # Set grid lines lighter
        ax1.grid(True, linestyle=':', linewidth=0.5, color='gray')
        ax2.grid(True, linestyle=':', linewidth=0.5, color='gray')

        # Enable cursor to display values on graphs
        mplcursors.cursor([line1, line2, line3, line4, line5])

        # Save the plot as an image or display it
        plt.tight_layout()  # Adjust layout to prevent clipping of labels
        # plt.savefig('graph.png')  # Save the plot as an image
        subfolder_path = os.path.join(main_folder_path, folder_name)
        os.makedirs(subfolder_path, exist_ok=True)
        plt.savefig(os.path.join(deeper_subfolder_path, 'graph.png'))  # Save the plot as an image in the specified directory
        plt.show()
        
        if 'localtime' not in data.columns:
            # Convert the 'timestamp' column from Unix milliseconds to datetime
            data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')

            # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
            data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)


        # Convert 'localtime' column to datetime
        data['localtime'] = pd.to_datetime(data['localtime'], format='%Y-%m-%d %H:%M:%S.%f')
        


    # Iterate through all subfolders in the main folder
    for subfolder in os.listdir(main_folder_path):
        subfolder_path = os.path.join(main_folder_path, subfolder)
        print("Subfolder name or Date of Ride=", subfolder)
        if os.path.isdir(subfolder_path):
            # Iterate through next level of subfolders
            for deeper_subfolder in os.listdir(subfolder_path):
                print("Ride=", deeper_subfolder)
                deeper_subfolder_path = os.path.join(subfolder_path, deeper_subfolder)
                if os.path.isdir(deeper_subfolder_path):
                    # Find the log file in the deepest subfolder
                    log_withoutanomaly_path = os.path.join(deeper_subfolder_path, 'log_withoutanomaly.csv')
                    log_file_path = os.path.join(deeper_subfolder_path, 'log_km.csv')
                    if not os.path.exists(log_withoutanomaly_path) and  os.path.exists(log_file_path):
                        data = pd.read_csv(log_file_path)
                    
                        # Convert 'timestamp' to datetime
                        data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms') + pd.Timedelta(hours=5, minutes=30)
            
                        plot_ghps(data, subfolder)
            
                        continue_removing = True
                        while continue_removing:
                            # Check if cropping is needed
                            crop = input("Do you want to remove anomalies? (yes/no): ")
                            if crop.lower() == "yes":
                                # Allow user to input start time and end time for the anomalies after seeing the graph
                                start_time_input = input("Enter start time of anomaly (format: DD-MM-YYYY HH:MM:SS): ")
                                end_time_input = input("Enter end time of anomaly (format: DD-MM-YYYY HH:MM:SS): ")

                                # Convert input strings to datetime
                                start_time = pd.to_datetime(start_time_input, format='%Y-%m-%d %H:%M:%S.%f')
                                end_time = pd.to_datetime(end_time_input, format='%Y-%m-%d %H:%M:%S.%f')

                                # Filter the data to exclude only the anomaly data between user-specified start and end times
                                data = data[(data['localtime'] < start_time) | (data['localtime'] > end_time)]

                                # Replot and save the normal data
                                plot_ghps(data, subfolder)  # Plotting the data without the anomalies
                            else:
                                continue_removing = False
                                # Define output file path for this Battery folder
                                output_file_path = os.path.join(deeper_subfolder_path, 'log_withoutanomaly.csv')

                                # Save the final filtered data to CSV in the same folder
                                data.to_csv(output_file_path, index=False)
                    else:
                        if os.path.exists(log_withoutanomaly_path):
                            print("Log without anomaly already exists for folder:", deeper_subfolder)
                        print("No log file found in folder:", deeper_subfolder)


    
def analysis(main_folder_path):
    from matplotlib.mlab import window_none
    import pandas as pd
    from math import radians, sin, cos, sqrt, atan2
    import sys
    import io
    import openpyxl
    import matplotlib.pyplot as plt
    from mpl_toolkits.basemap import Basemap
    from PIL import Image
    import os
    import matplotlib.dates as mdates
    import mplcursors  # Import mplcursors
    from contextlib import redirect_stdout
    from pptx.util import Inches, Pt  # Correcting the import statement
    from pptx import Presentation
    from pptx.util import Inches
    from docx import Document
    from docx.shared import Inches
    from openpyxl import load_workbook, Workbook
    from collections import defaultdict
    window_size =5
    
    # List to store DataFrames from each CSV file
    dfs = []
    def haversine(lat1, lon1, lat2, lon2):
        """
        Calculate the great-circle distance between two points
        on the Earth's surface using the Haversine formula.
        """
        # Convert latitude and longitude from degrees to radians
        lat1 = radians(lat1)
        lon1 = radians(lon1)
        lat2 = radians(lat2)
        lon2 = radians(lon2)
    
        # Radius of the Earth in kilometers
        R = 6371.0
    
        # Calculate differences in latitude and longitude
        dlat = lat2 - lat1
        dlon = lon2 - lon1
    
        # Calculate Haversine formula
        a = sin(dlat / 2)**2 + cos(lat1) * cos(lat2) * sin(dlon / 2)**2
        c = 2 * atan2(sqrt(a), sqrt(1 - a))
        distance = R * c
    
        return distance
    
    # Define a function to set current to zero if RPM is zero for 10 or more consecutive points
    def adjust_current(row):
        adjust_current.zero_count = getattr(adjust_current, 'zero_count', 0)
        if row['MotorSpeed_340920578'] == 0:
            adjust_current.zero_count += 1
        else:
            adjust_current.zero_count = 0
    
        if adjust_current.zero_count >= 10:
            return 0
        else:
            return row['PackCurr_6']
    
    
    
    def plot_ghps(data,Path):
        if 'localtime' not in data.columns:
            # Convert the 'timestamp' column from Unix milliseconds to datetime
            data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')
    
            # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
            data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)
    
        # data['localtime'] = pd.to_datetime(data['localtime'], format='%d/%m/%Y %H:%M:%S.%f', dayfirst=True)
        data['localtime'] = pd.to_datetime(data['localtime'])
        data.set_index('localtime', inplace=True)
        data['PackCurr_6'] = data.apply(adjust_current, axis=1)
    
        # Create a figure and axes for plotting
        fig, ax1 = plt.subplots(figsize=(10, 6))
    
        # Plot 'PackCurr_6' on primary y-axis
        line1, = ax1.plot(data.index, -data['PackCurr_6'], color='blue', label='PackCurr_6')
        ax1.set_ylabel('Pack Current (A)', color='blue')
        ax1.yaxis.set_label_coords(-0.1, 0.7)  # Adjust label position
    
        # Create secondary y-axis for 'MotorSpeed_340920578' (RPM)
        ax2 = ax1.twinx()
        line2, = ax2.plot(data.index, data['MotorSpeed_340920578'], color='green', label='Motor Speed')
        ax2.set_ylabel('Motor Speed (RPM)', color='green')
    
        # Add 'AC_Current_340920579' to primary y-axis
        line3, = ax1.plot(data.index, data['AC_Current_340920579'], color='red', label='AC Current')
    
        # Add 'AC_Voltage_340920580' scaled to 10x to the left side y-axis
        line4, = ax1.plot(data.index, data['AC_Voltage_340920580'] * 10, color='yellow', label='AC Voltage (x10)')
    
        # Add 'Throttle_408094978' to the left side y-axis
        line5, = ax1.plot(data.index, data['Throttle_408094978'], color='orange', label='Throttle (%)')
    
        # Add 'Throttle_408094978' to the left side y-axis
        line6, = ax1.plot(data.index, data['SOC_8'], color='black', label='SOC (%)')
    
    
        # Hide the y-axis label for 'AC_Current_340920579'
        ax1.get_yaxis().get_label().set_visible(False)
    
        # Set x-axis label and legend
        ax1.set_xlabel('Local Time')
        ax1.legend(loc='upper left')
        ax2.legend(loc='upper right')
    
        # Add a title to the plot
        plt.title('Battery Pack, Motor Data, and Throttle')
    
        # Format x-axis ticks as hours:minutes:seconds
        ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M:%S'))
    
        # Set grid lines lighter
        ax1.grid(True, linestyle=':', linewidth=0.5, color='gray')
        ax2.grid(True, linestyle=':', linewidth=0.5, color='gray')
    
        # Enable cursor to display values on graphs
        mplcursors.cursor([line1, line2, line3, line4, line5,line6])
    
        # Save the plot as an image or display it
        plt.tight_layout()  # Adjust layout to prevent clipping of labels
        # plt.savefig('graph.png')  # Save the plot as an image
    
        os.makedirs(Path, exist_ok=True)
        graph_path = os.path.join(Path, 'graph.png')
        plt.savefig(os.path.join(Path, 'graph.png'))  # Save the plot as an image in the specified directory
        # plt.show()
    
        return graph_path # Return the path of the saved graph image
    
    
    # def analysis_Energy(log_file, km_file):
    def analysis_Energy(log_file):
        dayfirst=True
        data = pd.read_csv(log_file)
    
        if 'localtime' not in data.columns:
        # Convert the 'timestamp' column from Unix milliseconds to datetime
            data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')
    
        # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
            data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)
    
    
        total_duration = 0
        total_distance = 0
        Wh_km = 0
        SOC_consumed = 0
    
        if (data['SOC_8'] > 90).any():
            # Maximum cell temperature calculation
            temp_columns_max = [f'Temp{i}_10' for i in range(1, 9)]
            max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
            max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
            print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
    
            # Minimum cell temperature calculation
            temp_columns_min = [f'Temp{i}_10' for i in range(1, 9)]
            min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
            min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
            print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
    
            #Difference between Maximum and Minimum cell Temperature
            CellTempDiff= max_cell_temp-min_cell_temp
            print("Temperature difference: ",CellTempDiff)
    
        else:
            print("SOC is not maximum!")
            # Maximum cell temperature calculation
            temp_columns_max = [f'Temp{i}_10' for i in range(1, 9)]
            max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
            max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
            print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
    
            # Minimum cell temperature calculation
            temp_columns_min = [f'Temp{i}_10' for i in range(1, 9)]
            min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
            min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
            print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
    
            #Difference between Maximum and Minimum cell Temperature
            CellTempDiff= max_cell_temp-min_cell_temp
            print("Temperature difference: ",CellTempDiff)
    
    
    
    
    
        # Check if 'localtime' column exists in data DataFrame
        if 'localtime' not in data.columns:
            print("Error: 'localtime' column not found in the DataFrame.")
            return None, None, None, None
    
        # Drop rows with missing values in 'SOCAh_8' column
        data.dropna(subset=['SOCAh_8'], inplace=True)
    
        # Convert the 'localtime' column to datetime s
        # data['localtime'] = pd.to_datetime(data['localtime'], format='%d/%m/%Y %H:%M:%S.%f',dayfirst=True)
        data['localtime'] = pd.to_datetime(data['localtime'])
    
        # Calculate the start time and end time
        start_time = data['localtime'].min()
        end_time = data['localtime'].max()
    
        # Calculate the start time and end time with formatting
        start_time_seconds = data['localtime'].min().strftime('%d/%m/%Y %H:%M:%S')
        end_time_seconds = data['localtime'].max().strftime('%d/%m/%Y %H:%M:%S')
    
        print(start_time,end_time)
    
        # Calculate the total time taken for the ride``
        total_duration = end_time - start_time
        total_hours = total_duration.seconds // 3600
        total_minutes = (total_duration.seconds % 3600) // 60
    
        print(f"Total time taken for the ride: {int(total_hours):02d}:{int(total_minutes):02d}")
    
        # Calculate the time difference between consecutive rows
        data['Time_Diff'] = data['localtime'].diff().dt.total_seconds().fillna(0)
    
        # Set the 'localtime' column as the index
        data.set_index('localtime', inplace=True)
    
        # Resample the data to have one-second intervals and fill missing values with previous ones
        data_resampled = data.resample('s').ffill()
    
        # Calculate the time difference between consecutive rows
        data_resampled['Time_Diff'] = data_resampled.index.to_series().diff().dt.total_seconds().fillna(0)
    
        # Calculate the actual Ampere-hours (Ah) using the trapezoidal rule for numerical integration
        actual_ah = abs((data_resampled['PackCurr_6'] * data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours
        print("Actual Ampere-hours (Ah): {:.2f}".format(actual_ah))
    
        # Calculate the actual Watt-hours (Wh) using the trapezoidal rule for numerical integration
        watt_h = abs((data_resampled['PackCurr_6'] * data_resampled['PackVol_6'] * data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours
        print("Actual Watt-hours (Wh):{:.2f}" .format(watt_h))
    
        ###########   starting and ending ah
        starting_soc_Ah = data['SOCAh_8'].iloc[-1]
        ending_soc_Ah = data['SOCAh_8'].iloc[0]
    
        print("Starting SoC (Ah):{:.2f}".format (starting_soc_Ah))
        print("Ending SoC (Ah):{:.2f}".format  (ending_soc_Ah))
    
    
    #Code for SOC_percentage(Starting and Ending SOC)
        starting_soc_percentage = data['SOC_8'].max()
        ending_soc_percentage = data['SOC_8'].min()
        print("Starting SOC:", starting_soc_percentage)
        print("Ending SOC:", ending_soc_percentage)
    
    
        # Initialize total distance covered
        total_distance = 0
        distance_per_mode = defaultdict(float)
    
        # Iterate over rows to compute distance covered between consecutive points
        for i in range(len(data) - 1):
        
            # Get latitude and longitude of consecutive points
            lat1 = data['latitude'].iloc[i]
            lon1 = data['longitude'].iloc[i]
            lat2 = data['latitude'].iloc[i + 1]
            lon2 = data['longitude'].iloc[i + 1]
        
    
            # Calculate distance between consecutive points
            distance = haversine(lat1, lon1, lat2, lon2)
    
            # Add distance to total distance covered
            total_distance += distance

              # Add distance to distance per mode
            mode = data['Mode_Ack_408094978'].iloc[i]
            distance_per_mode[mode] += distance
    
        print("Total distance covered (in kilometers):{:.2f}".format(total_distance))
    
        ##############   Wh/Km
        Wh_km = abs(watt_h / total_distance)
        print("WH/KM:{:.2f}". format (watt_h / total_distance))
    
        # Assuming 'data' is your DataFrame with 'SOC_8' column
        initial_soc = data['SOC_8'].iloc[-1]  # Initial SOC percentage
        final_soc = data['SOC_8'].iloc[0]   # Final SOC percentage
    
        # Calculate total SOC consumed
        total_soc_consumed =  abs(final_soc - initial_soc)
    
        print ("Total SOC consumed:{:.2f}".format (total_soc_consumed),"%")
    
    
        # Check if the mode remains constant or changes
        mode_values = data['Mode_Ack_408094978'].unique()
    
        if len(mode_values) == 1:
            # Mode remains constant throughout the log file
            mode = mode_values[0]
            if mode == 3:
                print("Mode is Custom mode.")
            elif mode == 2:
                print("Mode is Power mode.")
            elif mode == 1:
                print("Mode is Eco mode.")
        else:
            # Mode changes throughout the log file
            mode_counts = data['Mode_Ack_408094978'].value_counts(normalize=True) * 100
            for mode, percentage in mode_counts.items():
                if mode == 3:
                    print(f"Custom mode: {percentage:.2f}%")
                elif mode == 2:
                    print(f"Power mode: {percentage:.2f}%")
                elif mode == 1:
                    print(f"Eco mode: {percentage:.2f}%")
    
    
        ######################################
       # Function to calculate Wh/km for a given mode
        def calculate_wh_per_km(data, mode, distance):
            if distance < 1:
                return 0  # Return 0 if distance is less than 1 km
            mode_data = data[data['Mode_Ack_408094978'] == mode]
            if mode_data.empty:
                return 0  # Return 0 if there's no data for the given mode
            watt_h_mode = abs((mode_data['PackCurr_6'] * mode_data['PackVol_6'] * mode_data['Time_Diff']).sum()) / 3600
            return abs(watt_h_mode / distance)

        # Calculate Wh/km for each mode
        wh_per_km_CUSTOM_mode = calculate_wh_per_km(data_resampled, 3, distance_per_mode[3])
        print("Custom mode distance-------------->",distance_per_mode[3])
        wh_per_km_POWER_mode = calculate_wh_per_km(data_resampled, 2, distance_per_mode[2])
        print("POWER mode distance-------------->",distance_per_mode[2])
        wh_per_km_ECO_mode = calculate_wh_per_km(data_resampled, 1, distance_per_mode[1])
        print("ECO mode distance-------------->",distance_per_mode[1])

        # Calculate Wh/km for the entire ride
        watt_h = abs((data_resampled['PackCurr_6'] * data_resampled['PackVol_6'] * data_resampled['Time_Diff']).sum()) / 3600
        wh_per_km_total = abs(watt_h / total_distance)

        # Print the results
        print(f"Wh/km for Custom mode: {wh_per_km_CUSTOM_mode:.2f}")
        print(f"Wh/km for Power mode: {wh_per_km_POWER_mode:.2f}")
        print(f"Wh/km for Eco mode: {wh_per_km_ECO_mode:.2f}")
        print(f"Wh/km for the entire ride: {wh_per_km_total:.2f}")

        ######################################
    
        # Calculate power using PackCurr_6 and PackVol_6
        data_resampled['Power'] = data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
    
        # Find the peak power
        peak_power = data_resampled['Power'].min()
        print("Peak Power:", peak_power)
    
        # Calculate the average power
        average_power = abs(data_resampled['Power'].mean())
        print("Average Power:", average_power)
        altitude = None
        # Calculate altitude if 'Altitude' column exists in data_resampled
        if 'Altitude' in data_resampled.columns:
            altitude = data_resampled['Altitude'].max()  # Assuming you want to find the maximum altitude
            print("Maximum Altitude:", altitude)
        else:
            print("Altitude data not available.")
    
    
        # Calculate energy regenerated (in watt-hours)
        energy_regenerated = ((data_resampled[data_resampled['Power'] > 0]['Power']*data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
    
        # Calculate total energy consumed (in watt-hours)
        total_energy_consumed =  ((data_resampled[data_resampled['Power'] < 0]['Power']*data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
    
        print("total",total_energy_consumed)
        print("total energy regenerated",energy_regenerated)
    
        # Calculate regenerative effectiveness as a percentage
        if total_energy_consumed != 0:
            regenerative_effectiveness = abs(energy_regenerated / total_energy_consumed) * 100
            print("Regenerative Effectiveness (%):", regenerative_effectiveness)
        else:
            print("Total energy consumed is 0, cannot calculate regenerative effectiveness.")
    
        # Calculate idling time percentage (RPM was zero for more than 5 seconds)
        idling_time = ((data['MotorSpeed_340920578'] <= 0) | data['MotorSpeed_340920578'].isna()).sum()
        print("data length time",len(data))
        idling_percentage = (idling_time / len(data)) * 100
        print("Idling time percentage:", idling_percentage)
    
        # Calculate time spent in specific speed ranges
        speed_ranges = [(0, 10), (10, 20), (20, 30), (30, 40), (40, 50),(50, 60),(60,70),(70, 80),(80, 90)]
        speed_range_percentages = {}
    
        for range_ in speed_ranges:
            speed_range_time = ((data['MotorSpeed_340920578'] * 0.016 > range_[0]) & (data['MotorSpeed_340920578'] * 0.016 < range_[1])).sum()
            speed_range_percentage = (speed_range_time / len(data)) * 100
            speed_range_percentages[f"Time spent in {range_[0]}-{range_[1]} km/h"] = speed_range_percentage
            print(f"Time spent in {range_[0]}-{range_[1]} km/h: {speed_range_percentage:.2f}%")
    
            
        # Calculate power using PackCurr_6 and PackVol_6
        data_resampled['Power'] = -data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
    
        # Find the peak power
        peak_power = data_resampled['Power'].max()
    
        # Get the maximum cell voltage
        max_cell_voltage = data_resampled['MaxCellVol_5'].max()
    
        # Find the index where the maximum voltage occurs
        max_index = data_resampled['MaxCellVol_5'].idxmax()
    
        # Retrieve the corresponding cell ID using the index
        max_cell_id = data_resampled['MaxVoltId_5'].loc[max_index]
    
        # Get the minimum cell voltage
        min_cell_voltage = data_resampled['MinCellVol_5'].min()
    
        # Find the index where the minimum voltage occurs
        min_index = data_resampled['MinCellVol_5'].idxmin()
    
        # Retrieve the corresponding cell ID using the index
        min_cell_id = data_resampled['MinVoltId_5'].loc[min_index]
    
        voltage_difference = max_cell_voltage - min_cell_voltage
    
        # Get the maximum temperature
        max_temp = data_resampled['MaxTemp_7'].max()
    
        # Find the index where the maximum temperature occurs
        max_temp_index = data_resampled['MaxTemp_7'].idxmax()
    
        # Retrieve the corresponding temperature ID using the index
        max_temp_id = data_resampled['MaxTempId_7'].loc[max_temp_index]
    
        # Get the minimum temperature
        min_temp = data_resampled['MinTemp_7'].min()
    
        # Find the index where the minimum temperature occurs
        min_temp_index = data_resampled['MinTemp_7'].idxmin()
    
        # Retrieve the corresponding temperature ID using the index
        min_temp_id = data_resampled['MinTempId_7'].loc[min_temp_index]
    
        # Calculate the difference in temperature
        temp_difference = max_temp - min_temp
    
        # Get the maximum temperature of FetTemp_8
        max_fet_temp = data_resampled['FetTemp_8'].max()
    
        # Get the maximum temperature of AfeTemp_12
        max_afe_temp = data_resampled['AfeTemp_12'].max()
    
        # Get the maximum temperature of PcbTemp_12
        max_pcb_temp = data_resampled['PcbTemp_12'].max()
    
        # Get the maximum temperature of MCU_Temperature_408094979
        max_mcu_temp = data_resampled['MCU_Temperature_408094979'].max()
    
        # Check for abnormal motor temperature at high RPMs
        max_motor_temp = data_resampled['Motor_Temperature_408094979'].max()
    
        # Find the battery voltage
        batteryVoltage = (data_resampled['BatteryVoltage_340920578'].max()) * 10
        print( "Battery Voltage", batteryVoltage )
    
        # Check for abnormal motor temperature at high RPMs for at least 15 seconds
        abnormal_motor_temp = (data_resampled['Motor_Temperature_408094979'] < 10) & (data_resampled['MotorSpeed_340920578'] > 3500)
        abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
    
        # Check if abnormal condition persists for at least 15 seconds
        abnormal_motor_temp_detected = (abnormal_motor_temp_mask >= 120).any()
    
        #For battery Analysis
        cycleCount= data_resampled['CycleCount_7'].max()

        SOH = data_resampled['SOH_8'].max()
    
    
    
        ByteBeamId= data['id'].iloc[0]
    
        max_continuous_duration = 0
    
        for speed in range(int(data_resampled['MotorSpeed_340920578'].min()), int(data_resampled['MotorSpeed_340920578'].max()) + 1):
            lower_bound = speed - window_size
            upper_bound = speed + window_size
        
    
            within_window = data_resampled[(data_resampled['MotorSpeed_340920578'] >= lower_bound) & (data_resampled['MotorSpeed_340920578'] <= upper_bound)].copy()
            within_window.loc[:, 'Group'] = (within_window['MotorSpeed_340920578'].diff() > window_size).cumsum()
            continuous_durations = within_window.groupby('Group')['Time_Diff'].sum()
        
            current_max_duration = continuous_durations.max() if not continuous_durations.empty else 0
    
            if current_max_duration > max_continuous_duration:
                print("max_continuous_duration------->",max_continuous_duration)
                max_continuous_duration = current_max_duration
                cruising_rpm = speed
                cruising_speed=speed*0.01606

                if cruising_speed >1:
                    cruise_speed=cruising_speed
                else :
                    cruise_speed =0
    
    
            # Find the maximum value in the 'MotorSpeed_340920578' column
        Max_motor_rpm = data_resampled['MotorSpeed_340920578'].max()
    
        # Print the maximum motor speed in RPM
        print("The maximum motor speed in RPM is:", Max_motor_rpm)
    
        # Convert the maximum motor speed to speed using the given factor
        peak_speed = Max_motor_rpm * 0.01606
    
        # Print the maximum speed
        print("The maximum speed is:", peak_speed)
    
        total_energy_kwh = actual_ah * batteryVoltage / 1000
        print("Total energy charged in kWh: {:.2f}".format(total_energy_kwh))
    
        total_energy_kw = total_energy_kwh / total_duration.seconds / 3600
        print("Electricity consumption units in kW", (total_energy_kw))

        print("wh/km----------->",wh_per_km_total)
    
        # Add these variables and logic to ppt_data
    
        ppt_data = {
            "Date and Time": str(start_time_seconds) + " to " + str(end_time_seconds),
            "Byte Beam ID": ByteBeamId,
            "Total time taken for the ride": total_duration,
            "Actual Ampere-hours (Ah)": actual_ah,
            "Actual Watt-hours (Wh)- Calculated_UsingFormala- P= VI": watt_h,
            "Starting SoC (Ah)": starting_soc_Ah,
            "Ending SoC (Ah)": ending_soc_Ah,
            "Starting SoC (%)": starting_soc_percentage,
            "Ending SoC (%)": ending_soc_percentage,
            "SOH": SOH,
            "Total distance covered (km)": total_distance,
            "Total energy consumption(WH/KM)- ENTIRE RIDE": watt_h / total_distance,
            "Total SOC consumed(%)":starting_soc_percentage- ending_soc_percentage,
            "Mode": "",
            "Wh/km in CUSTOM mode": wh_per_km_CUSTOM_mode,
            "Distance travelled in Custom mode":distance_per_mode[3],
            "Wh/km in POWER mode": wh_per_km_POWER_mode,
            "Distance travelled in POWER mode":distance_per_mode[2],
            "Wh/km in ECO mode": wh_per_km_ECO_mode,
            "Distance travelled in ECO mode":distance_per_mode[1],
            "Peak Power(kW)": peak_power,
            "Average Power(kW)": average_power,
            "Total Energy Regenerated(kWh)": energy_regenerated,
            "Regenerative Effectiveness(%)": regenerative_effectiveness,
            "Highest Cell Voltage(V)": max_cell_voltage,
            "Lowest Cell Voltage(V)": min_cell_voltage,
            "Difference in Cell Voltage(V)": voltage_difference,
            "Minimum Temperature(C)": min_temp,
            "Maximum Temperature(C)": max_temp,
            "Difference in Temperature(C)": max_temp- min_temp,
            "Maximum Fet Temperature-BMS(C)": max_fet_temp,
            "Maximum Afe Temperature-BMS(C)": max_afe_temp,
            "Maximum PCB Temperature-BMS(C)": max_pcb_temp,
            "Maximum MCU Temperature(C)": max_mcu_temp,
            "Maximum Motor Temperature(C)": max_motor_temp,
            "Abnormal Motor Temperature Detected(C)": abnormal_motor_temp_detected,
            "highest cell temp(C)": max_cell_temp,
            "lowest cell temp(C)": min_cell_temp,
            "Difference between Highest and Lowest Cell Temperature at 100% SOC(C)": CellTempDiff,
            "Battery Voltage(V)": batteryVoltage,
            "Total energy charged(kWh)- Calculated_Using_BatteryData": total_energy_kwh,
            "Electricity consumption units(kW)": total_energy_kw,
            "Cycle Count of battery": cycleCount,
            "Cruising Speed (Rpm)": cruising_rpm,
            "cruising_speed (km/hr)":cruise_speed,
            "Maximum Motor speed (RPM)":Max_motor_rpm,
            "Peak speed (Km/hr)":peak_speed
            }
    
        mode_values = data_resampled['Mode_Ack_408094978'].unique()
        if len(mode_values) == 1:
            mode = mode_values[0]
            if mode == 3:
                ppt_data["Mode"] = "Custom mode"
            elif mode == 2:
                ppt_data["Mode"] = "Power mode"
            elif mode == 1:
                ppt_data["Mode"] = "Eco mode"
        else:
            # Mode changes throughout the log file
            mode_counts = data_resampled['Mode_Ack_408094978'].value_counts(normalize=True) * 100
            mode_strings = []  # Initialize list to store mode strings
            for mode, percentage in mode_counts.items():
                if mode == 3:
                    mode_strings.append(f"Custom mode\n{percentage:.2f}%")
                elif mode == 2:
                    mode_strings.append(f"Power mode\n{percentage:.2f}%")
                elif mode == 1:
                    mode_strings.append(f"Eco mode\n{percentage:.2f}%")
            ppt_data["Mode"] = "\n".join(mode_strings)
    
        # Add calculated parameters to ppt_data
        ppt_data["Idling time percentage"] = idling_percentage
        ppt_data.update(speed_range_percentages)
    
        # Calculate power using PackCurr_6 and PackVol_6
        data_resampled['Power'] = -data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
    
        # Find the peak power
        peak_power = data_resampled['Power'].max()
        print("Peak Power:", peak_power)
    
        # Get the maximum cell voltage
        max_cell_voltage = data_resampled['MaxCellVol_5'].max()
    
        # Find the index where the maximum voltage occurs
        max_index = data_resampled['MaxCellVol_5'].idxmax()
    
        # Retrieve the corresponding cell ID using the index
        max_cell_id = data_resampled['MaxVoltId_5'].loc[max_index]
    
            # Get the minimum cell voltage
        min_cell_voltage = data_resampled['MinCellVol_5'].min()
    
        # Find the index where the minimum voltage occurs
        min_index = data_resampled['MinCellVol_5'].idxmin()
    
        # Retrieve the corresponding cell ID using the index
        min_cell_id = data_resampled['MinVoltId_5'].loc[min_index]
    
        voltage_difference = max_cell_voltage - min_cell_voltage
    
    
        print("Lowest Cell Voltage:", min_cell_voltage, "V, Cell ID:", min_cell_id)
        print("Highest Cell Voltage:", max_cell_voltage, "V, Cell ID:", max_cell_id)
        print("Difference in Cell Voltage:", voltage_difference, "V")
    
        # Get the maximum temperature
        max_temp = data_resampled['MaxTemp_7'].max()
    
        # Find the index where the maximum temperature occurs
        max_temp_index = data_resampled['MaxTemp_7'].idxmax()
    
        # Retrieve the corresponding temperature ID using the index
        max_temp_id = data_resampled['MaxTempId_7'].loc[max_temp_index]
    
    
        # Get the minimum temperature
        min_temp = data_resampled['MinTemp_7'].min()
    
        # Find the index where the minimum temperature occurs
        min_temp_index = data_resampled['MinTemp_7'].idxmin()
    
        # Retrieve the corresponding temperature ID using the index
        min_temp_id = data_resampled['MinTempId_7'].loc[min_temp_index]
        # Calculate the difference in temperature
        temp_difference = max_temp - min_temp
        print("temp_difference: ",temp_difference)
        print("Total distance: ",total_distance)
    
        print("Maximum Temperature:", max_temp, "C, Temperature ID:", max_temp_id)
        print("Minimum Temperature:", min_temp, "C, Temperature ID:", min_temp_id)
        print("Difference in Temperature:", temp_difference, "C")
    
        # Get the maximum temperature of FetTemp_8
        max_fet_temp = data_resampled['FetTemp_8'].max()
        print("Maximum Fet Temperature:", max_fet_temp, "C")
    
        # Get the maximum temperature of AfeTemp_12
        max_afe_temp = data_resampled['AfeTemp_12'].max()
        print("Maximum Afe Temperature:", max_afe_temp, "C")
    
        # Get the maximum temperature of PcbTemp_12
        max_pcb_temp = data_resampled['PcbTemp_12'].max()
        print("Maximum PCB Temperature:", max_pcb_temp, "C")
    
        # Get the maximum temperature of MCU_Temperature_408094979
        max_mcu_temp = data_resampled['MCU_Temperature_408094979'].max()
        print("Maximum MCU Temperature:", max_mcu_temp, "C")
    
        # Check for abnormal motor temperature at high RPMs
        max_motor_temp = data_resampled['Motor_Temperature_408094979'].max()
    
    
        print("Maximum Motor Temperature:", max_motor_temp, "C")
    
    
        # Check for abnormal motor temperature at high RPMs for at least 15 seconds
        abnormal_motor_temp = (data_resampled['Motor_Temperature_408094979'] < 10) & (data_resampled['MotorSpeed_340920578'] > 3500)
    
        # Convert to a binary mask indicating consecutive occurrences
        abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
    
        return total_duration, total_distance, Wh_km, total_soc_consumed,ppt_data
    
    
    
    # def capture_analysis_output(log_file, km_file, folder_path):
    def capture_analysis_output(log_file,folder_path):
        try:
            analysis_output = io.StringIO()
            output_file = "analysis_results.docx"
    
            with redirect_stdout(analysis_output):
                #total_duration, total_distance, Wh_km, total_soc_consumed, ppt_data = analysis_Energy(log_file, km_file)
                analysis_output = analysis_output.getvalue()
    
            # Extract folder name from folder_path
            folder_name = os.path.basename(folder_path)
    
            # Create a new PowerPoint presentation
            prs = Presentation()
    
            # Add title slide with 'Selawik' style
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = f"Analysis Results from Folder - {folder_name}"
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.size = Pt(36)  # Corrected to Pt
            title.text_frame.paragraphs[0].font.name = 'Selawik'
    
            rows = len(ppt_data) + 1
            cols = 2
            table_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(table_slide_layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = "Analysis Results:"
    
            # Centering the title horizontally
            title_shape.left = Inches(0.5)
            title_shape.top = Inches(0.5)  # Adjust as needed
    
            # Setting the font size of the title
            title_shape.text_frame.paragraphs
    
            # Define maximum number of rows per slide
            max_rows_per_slide = 13
            # Add some space between title and table
            title_shape.top = Inches(0.5)
    
            table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
            table.columns[0].width = Inches(4)
            table.columns[1].width = Inches(4)
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"
    
            # Initialize row index
            row_index = 1
    
            # Iterate over data and populate the table
            for key, value in ppt_data.items():
                # Check if current slide has reached maximum rows
                if row_index > max_rows_per_slide:
                    # Add a new slide
                    slide = prs.slides.add_slide(table_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    title_shape.text = "Analysis Results:"
    
                    # Add a new table to the new slide
                    table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
                    table.columns[0].width = Inches(4)
                    table.columns[1].width = Inches(4)
                    table.cell(0, 0).text = "Metric"
                    table.cell(0, 1).text = "Value"
    
                    # Reset row index for the new slide
                    row_index = 1
    
                # Populate the table
                table.cell(row_index, 0).text = key
                table.cell(row_index, 1).text = str(value)
    
                # Increment row index
                row_index += 1
    
            graph_path= plot_ghps(data,subfolder_path)
    
            # Add image slide with title and properly scaled image
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
    
            # Remove the unwanted title placeholder
            for shape in slide.shapes:
                if shape.is_placeholder:
                    slide.shapes._spTree.remove(shape._element)
    
            # Add the title
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), prs.slide_width - Inches(2), Inches(1))
            title_shape.text = "Graph Analysis"
    
            # Add the image and adjust its position and size
            graph_width = prs.slide_width - Inches(1)
            graph_height = prs.slide_height - Inches(2)
            left = (prs.slide_width - graph_width) / 2
            top = (prs.slide_height - graph_height) / 2 + Inches(1)
            pic = slide.shapes.add_picture(graph_path, left, top, width=graph_width, height=graph_height)
    
            # Save the presentation
            output_file_name = f"{folder_path}/analysis_{folder_name}.pptx"
            prs.save(output_file_name)
            print("PPT generated!")
    
            # Create a new Excel workbook
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Analysis Results"
    
            # Populate Excel sheet with analysis data
            for i, (key, value) in enumerate(ppt_data.items(), start=1):
                ws.cell(row=i, column=1, value=key)
                ws.cell(row=i, column=2, value=value)
    
            # Save the Excel workbook
            excel_output_file = f"{folder_path}/analysis_{folder_name}.xlsx"
            wb.save(excel_output_file)
            print("Excel generated!")
    
        except Exception as e:
            print("Error:", e)
    
    
    
    
    
    
    
    # Initialize variables to store file paths
    log_file = None
    
    
    
    
    def mergeExcel(main_folder_path):
        def prepare_sheet_in_memory(file_path):
            workbook = load_workbook(filename=file_path)
            sheet = workbook.active
            file_name = os.path.basename(file_path)
            sheet.insert_rows(1)
            sheet['A1'] = 'file name'
            sheet['B1'] = file_name
            return sheet, file_name
    
        def sheet_to_dict(sheet):
            data_dict = {}
            for row in sheet.iter_rows(min_row=2, values_only=True):
                key, *values = row
                data_dict[key] = values
            return data_dict
    
        def merge_dicts(dict1, dict2):
            for key, values in dict2.items():
                if key in dict1:
                    # Assuming that values is a list of values
                    dict1[key].extend(values)
                else:
                    dict1[key] = values
            return dict1
    
        def process_directory(directory):
            merged_data = {}
            file_names = []
            for root, dirs, files in os.walk(directory):
                for name in dirs:
                    # if name.startswith("B"):
                        subdir_path = os.path.join(root, name)
                        for file_name in os.listdir(subdir_path):
                            if file_name.endswith(".xlsx"):
                                file_path = os.path.join(subdir_path, file_name)
                                sheet, extracted_file_name = prepare_sheet_in_memory(file_path)
                                data_dict = sheet_to_dict(sheet)
                                merged_data = merge_dicts(merged_data, data_dict)
                                if extracted_file_name not in file_names:
                                    file_names.append(extracted_file_name)
            return merged_data, file_names
    
        def main(main_folder_path):
            directory = main_folder_path
            merged_data, file_names = process_directory(directory)
    
            merged_workbook = Workbook()
            merged_sheet = merged_workbook.active
    
            # # Use the extracted file names for headers
            # headers = ['File name'] + file_names
            # merged_sheet.append(headers)
    
            if merged_data:
                for key, values in merged_data.items():
                    merged_sheet.append([key] + values)
            else:
                print("No data found in merged_data")
    
            merged_file_path = os.path.join(directory, 'ANALYSIS.xlsx')
            merged_workbook.save(filename=merged_file_path)
            print("Merged Excel file is ready")
    
        if __name__ == '__main__':
            main(main_folder_path)
    
    
    
    
    # Iterate through each subfolder in the main folder path
    for root, dirs, files in os.walk(main_folder_path):
        for dir_name in dirs:
            subfolder_path = os.path.join(root, dir_name)
            log_file_without_anomaly = None
            log_file_with_anomaly = None

            # Check within the directory for specific files
            for file in os.listdir(subfolder_path):
                if file == 'log_withoutanomaly.csv':
                    log_file_without_anomaly = os.path.join(subfolder_path, file)
                elif file == 'log_km.csv':
                    log_file_with_anomaly = os.path.join(subfolder_path, file)

            # Choose the appropriate file to process
            log_file_to_use = log_file_without_anomaly if log_file_without_anomaly else log_file_with_anomaly

            if log_file_to_use:
                print(f"Processing file: {log_file_to_use}")  # Debug print to verify files being processed
                try:
                    data = pd.read_csv(log_file_to_use)  # Read the data
                    total_duration, total_distance, Wh_km, SOC_consumed, ppt_data = analysis_Energy(log_file_to_use)
                    capture_analysis_output(log_file_to_use, os.path.dirname(log_file_to_use))

                    if 'data' in locals():
                        graph_path = plot_ghps(data, os.path.dirname(log_file_to_use))
                    else:
                        print("Data not available for plotting.")
                except Exception as e:
                    print(f"Failed to process {log_file_to_use}: {str(e)}")
            else:
                print(f"No suitable log file found in subfolder: {subfolder_path}")

    mergeExcel(main_folder_path)



Merge_log_km(main_folder_path)
crop_data(main_folder_path)
analysis(main_folder_path)
