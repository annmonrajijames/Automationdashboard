def Influx_LXS_input(input_folder_path):    
    import sys   
    from matplotlib.mlab import window_none
    import pandas as pd
    from math import radians, sin, cos, sqrt, atan2
    import io
    import openpyxl
    import matplotlib.pyplot as plt
    from PIL import Image
    import os
    import matplotlib.dates as mdates
    import mplcursors  # Import mplcursors
    from contextlib import redirect_stdout
    from pptx.util import Inches, Pt  # Correcting the import statement
    from pptx import Presentation
    from pptx.util import Inches
    from docx import Document
    from docx.shared import Inches
    from openpyxl import load_workbook, Workbook
    from openpyxl.drawing.image import Image
    from matplotlib.widgets import CheckButtons
    from collections import defaultdict
    import plotly.graph_objs as go
    from plotly.subplots import make_subplots
    window_size =5  

    def process_data(file_path):
        # Read the first line to get the date and time from the header
        with open(file_path, 'r') as file:
            first_line = file.readline().strip()
        date_time = first_line.split(':', 1)[1].strip() if ':' in first_line else first_line

        # Load the data using a chunk size
        chunk_size = 50000  # Adjust chunk size based on your system's memory
        chunks = pd.read_csv(file_path, skiprows=1, chunksize=chunk_size)
        
        # Check if the "Time" column exists in the first chunk
        first_chunk = next(chunks, None)
        if first_chunk is None or "Time" not in first_chunk.columns:
            return file_path  # If no "Time" column, return without processing
        
        # If "Time" column exists, reinitialize chunks iterator and proceed with processing
        chunks = pd.read_csv(file_path, skiprows=1, chunksize=chunk_size)
        data_frames = []
        
        for chunk in chunks:
            # Insert the date and time as a new column
            chunk.insert(0, 'Creation Time', date_time)
            data_frames.append(chunk)

        # Concatenate all chunks into one DataFrame
        data = pd.concat(data_frames, ignore_index=True)

        # Drop the first three rows which were initially 2nd, 3rd, and 4th in the original file
        data = data.drop([0, 1, 2])

        # Save the processed data to the same file, overwriting the original
        data.to_csv(file_path, index=False)
        
        return file_path

    def process_files_in_directory(root_directory):
        # Walk through all directories starting from the root
        for dirpath, dirnames, filenames in os.walk(root_directory):
            for filename in filenames:
                if filename == 'log.csv':  # Check for the specific filename 'log.csv'
                    file_path = os.path.join(dirpath, filename)
                    processed_file_path = process_data(file_path)
                    print(f"Processed file saved to: {processed_file_path}")

    # List to store DataFrames from each CSV file
    dfs = []
    def haversine(lat1, lon1, lat2, lon2):
        """
        Calculate the great-circle distance between two points
        on the Earth's surface using the Haversine formula.
        """
        # Convert latitude and longitude from degrees to radians
        lat1 = radians(lat1)
        lon1 = radians(lon1)
        lat2 = radians(lat2)
        lon2 = radians(lon2)
    
        # Radius of the Earth in kilometers
        R = 6371.0
    
        # Calculate differences in latitude and longitude
        dlat = lat2 - lat1
        dlon = lon2 - lon1
    
        # Calculate Haversine formula
        a = sin(dlat / 2)**2 + cos(lat1) * cos(lat2) * sin(dlon / 2)**2
        c = 2 * atan2(sqrt(a), sqrt(1 - a))
        distance = R * c
    
        return distance
    
    # Define a function to set current to zero if RPM is zero for 10 or more consecutive points
    def adjust_current(row):
        adjust_current.zero_count = getattr(adjust_current, 'zero_count', 0)
        if row['RPM'] == 0:
            adjust_current.zero_count += 1
        else:
            adjust_current.zero_count = 0
    
        if adjust_current.zero_count >= 10:
            return 0
        else:
            return row['Current_value']
    
    
    
    def plot_ghps(data,Path,maxCellTemp):
    
    
    
        speed = data['RPM'] * 0.0836
    
        data.set_index('DATETIME', inplace=True)  # Setting DATETIME as index
    
    
    
    
        fig = make_subplots(specs=[[{"secondary_y": True}]])
        fig.add_trace(go.Scatter(x=data.index, y=-data['Current_value'], name='Pack Current', line=dict(color='blue')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=data['RPM'], name='Motor Speed[RPM]', line=dict(color='green')), secondary_y=True)
        # fig.add_trace(go.Scatter(x=data.index, y=data['AC_Current [SA: 03]'], name='AC Current', line=dict(color='red')), secondary_y=False)
        # fig.add_trace(go.Scatter(x=data.index, y=data['AC_Voltage [SA: 04]'] * 10, name='AC Voltage (x10)', line=dict(color='yellow')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=data['Throttle [SA: 02]'], name='Throttle (%)', line=dict(color='Black')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=data['SOC [SA: 08]'], name='SOC (%)', line=dict(color='Red')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=speed, name='Speed[Km/hr]', line=dict(color='grey')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=data['voltage_value'], name='PackVoltage', line=dict(color='Green')), secondary_y=False)
        fig.add_trace(go.Scatter(x=data.index, y=data['ALTITUDE'], name='ALTITUDE', line=dict(color='red')), secondary_y=True)
        fig.add_trace(go.Scatter(x=data.index, y=data[maxCellTemp], name='MaximumCellTemperature', line=dict(color='orange')), secondary_y=True)
        fig.add_trace(go.Scatter(x=data.index, y=data['FetTemp [SA: 08]'], name='BMS temperature (FET)', line=dict(color='orange')), secondary_y=True)
        fig.add_trace(go.Scatter(x=data.index, y=data['MCS_Temp'], name='MCU temperature', line=dict(color='orange')), secondary_y=True)
        fig.add_trace(go.Scatter(x=data.index, y=data['Motor_Temp'], name='Motor Temperature', line=dict(color='orange')), secondary_y=True)
        # fig.add_trace(go.Scatter(x=data.index, y=data['Power'], name='DC Power', line=dict(color='Red')), secondary_y=True)
        fig.add_trace(go.Scatter(x=data.index, y=data['DeltaCellVoltage'], name='DeltaCellVoltage', line=dict(color='Purple')), secondary_y=True)
    
        fig.update_layout(title='Battery Pack, Motor Data, and Throttle',
                        xaxis_title='Local localtime',
                        yaxis_title='Pack Current (A)',
                        yaxis2_title='Motor Speed (RPM)')
    
        fig.update_xaxes(tickformat='%H:%M:%S')
    
        # Save the plot as an HTML file
        os.makedirs(Path, exist_ok=True)
        graph_path = os.path.join(Path, 'Log.html')
        fig.write_html(graph_path)
    
    
    # Resetting index to default integer index
        data.reset_index(inplace=True)
    
        # data.set_index('DATETIME', inplace=False)
    
    
        return graph_path  # Return the path of the saved graph image
    
    # Example usage:
    # data = pd.read_csv('your_data.csv')
    # Path = 'path/to/save'
    # plot_ghps(data, Path)
    
    
    
    # def analysis_Energy(log_file, km_file):
    def analysis_Energy(data,subfolder_path):
        dayfirst=True
    
        # data = pd.read_csv(log_file)
        # Remove duplicates based on the "DATETIME" column and keep the first occurrence
        # data = data.drop_duplicates(subset=['DATETIME'], keep='first')
    
    
    #Uncomment the following two lines if DateTime is already in IST.
        # data['DATETIME'] = pd.to_datetime(data['DATETIME'], unit='s', origin='unix').dt.strftime('%Y-%m-%d %H:%M:%S.%f').str[:-3]
        # data['DATETIME'] = pd.to_datetime(data['DATETIME'])
    
    
    #uncomment and run if ride contains only one ride(completely) and in that case dont need to run the split battery!
    
        # # Convert Unix epoch time to datetime, assuming the original timezone is UTC
        # data['DATETIME'] = pd.to_datetime(data['DATETIME'], unit='s', origin='unix', utc=True)
        # # Convert to your desired timezone (e.g., 'Asia/Kolkata')
        # data['DATETIME'] = data['DATETIME'].dt.tz_convert('Asia/Kolkata')  # converting to IST
        # # Format the datetime as string, including milliseconds
        # data['DATETIME'] = data['DATETIME'].dt.strftime('%Y-%m-%d %H:%M:%S.%f').str[:-3]  # Converting to string
        # # If you need the datetime back as pandas datetime type without timezone info
        # data['DATETIME'] = pd.to_datetime(data['DATETIME'])
    
        # data.set_index('DATETIME', inplace=True)
    
        total_duration = 0
        total_distance = 0
        Wh_km = 0
        mode_values = 0
    
    ###Commented the following because cell temperature values are not available in the LXS dbc.
        # if (data['SOC_value'] > 90).any():
        #     # Maximum cell temperature calculation
        #     temp_columns_max = [f'Temp{i} [SA: 0A]' for i in range(1, 9)]
        #     max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        #     max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
        #     print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
    
        #     # Minimum cell temperature calculation
        #     temp_columns_min = [f'Temp{i} [SA: 0A]' for i in range(1, 9)]
        #     min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        #     min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
        #     print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
    
        #     #Difference between Maximum and Minimum cell Temperature
        #     CellTempDiff= max_cell_temp-min_cell_temp
        #     print("Temperature difference: ",CellTempDiff)
    
        # else:
        #     print("SOC is not maximum!")
        #      # Maximum cell temperature calculation
        #     temp_columns_max = [f'Temp{i} [SA: 0A]' for i in range(1, 9)]
        #     max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        #     max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
        #     print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
    
        #     # Minimum cell temperature calculation
        #     temp_columns_min = [f'Temp{i} [SA: 0A]' for i in range(1, 9)]
        #     min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        #     min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
        #     print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
    
        #     #Difference between Maximum and Minimum cell Temperature
        #     CellTempDiff= max_cell_temp-min_cell_temp
        #     print("Temperature difference: ",CellTempDiff)
    
    
        # Define the temperature columns
        # temp_columns = [f'Temp{i} [SA: 0A]' for i in range(1, 9)]
    
        # Find the maximum temperature value among columns
        # max_values = data[temp_columns].max()
    
        # Determine which column(s) have the maximum temperature
        # max_columns_multiple = max_values[max_values == max_values.max()].index
    
        # If there are multiple columns with the same maximum temperature,
        # choose the column with the highest occurrence
        # if len(max_columns_multiple) > 1:
        #     max_occurrence_counts = data[max_columns_multiple].apply(lambda x: x.value_counts().get(max_values.max(), 0))
        #     max_column = max_occurrence_counts.idxmax()
        # else:
        #     max_column = max_columns_multiple[0]
    
        # Determine the maximum temperature occurrence in the selected column
        # max_occurrence = data[max_column].value_counts().idxmax()
    
        # Print the results
        # print(f"Overall maximum value of cell temperature: {max_values.max()}")
        # # print(f"Column(s) with maximum temperature: {max_columns_multiple}")
        # print(f"Column with the highest cell temperature: {max_column}")
        # # print(f"Maximum occurrence temperature: {max_occurrence}")
    
    
    
        ##################
        data['Power'] = data['Current_value'] * data['voltage_value']
    
    
    #####Commented because cell voltage not available in the LXS dbc
        # # Specify the columns of interest
        # columns_of_interest = [
        #     'CellVol01 [SA: 01]', 'CellVol02 [SA: 01]', 'CellVol03 [SA: 01]', 'CellVol04 [SA: 01]',
        #     'CellVol05 [SA: 02]', 'CellVol06 [SA: 02]', 'CellVol07 [SA: 02]', 'CellVol08 [SA: 02]',
        #     'CellVol09 [SA: 03]', 'CellVol10 [SA: 03]', 'CellVol11 [SA: 03]', 'CellVol12 [SA: 03]',
        #     'CellVol13 [SA: 04]', 'CellVol14 [SA: 04]', 'CellVol15 [SA: 04]', 'CellVol16 [SA: 04]'
        # ]
    
        # # Create an empty list to store computed differences
        # differences = []
    
        # # Iterate through each row and compute max, min, and their difference for each row
        # for index, row in data[columns_of_interest].iterrows():
        #     max_value = row.max()
        #     min_value = row.min()
        #     difference = max_value - min_value
        
        #     differences.append(difference)  # Append the computed difference to the list
    
        # # Add the differences list as a new column 'CellDifference' in the DataFrame
        # data['DeltaCellVoltage'] = differences
    
    
        # plot_ghps(data,subfolder_path,max_column)
    
        # Drop rows with missing values in 'SOCAh [SA: 08]' column
        data.dropna(subset=['SOC_value'], inplace=True)
    
    
        # # Calculate the start localtime and end localtime
        # start_localtime = data['DATETIME'].min()
        # end_localtime = data['DATETIME'].max()
    
        # # Calculate the start localtime and end localtime with formatting
        # start_localtime_seconds = data['DATETIME'].min().strflocaltime('%d/%m/%Y %H:%M:%S')
        # end_localtime_seconds = data['DATETIME'].max().strflocaltime('%d/%m/%Y %H:%M:%S')
        # data['DATETIME'] = data['DATETIME'].astype(str)
        # data['DATETIME'] = pd.to_datetime(data['DATETIME'].str.split('.').str[0], format='%Y-%m-%d  %H:%M:%S')
        
        # Convert the DATETIME column to datetime objects
        # Ensure the DATETIME column is read as integers
        data['DATETIME'] = pd.to_numeric(data['DATETIME'], errors='coerce')
    
        # Check for any NaN values which indicate parsing issues
        # print(data['DATETIME'].isnull().sum())
    
        # Drop or handle NaN values
        data = data.dropna(subset=['DATETIME'])
    
        # Convert the Unix timestamps to datetime
        data['DATETIME'] = pd.to_datetime(data['DATETIME'], unit='s')
    
        # Print the converted DATETIME column
        data['DATETIME'] = pd.to_datetime(data['DATETIME'])
        # Find the minimum and maximum datetime
        start_localtime = data['DATETIME'].min()
        end_localtime = data['DATETIME'].max()
    
        print("Start time of the ride: ",start_localtime)
        print("End time of the ride:",end_localtime)
    
        # # Format the start and end times
        # start_localtime_formatted = start_localtime.strftime('%d/%m/%Y %H:%M:%S')
        # end_localtime_formatted = end_localtime.strftime('%d/%m/%Y %H:%M:%S')
    
        # print(f"Start Local Time: {start_localtime_formatted}")
        # print(f"End Local Time: {end_localtime_formatted}")
        # ############
    
    
        # Calculate the total localtime taken for the ride``
        total_duration = end_localtime - start_localtime
        total_hours = total_duration.seconds // 3600
        total_minutes = (total_duration.seconds % 3600) // 60
    
        print(f"Total time taken for the ride: {int(total_hours):02d}:{int(total_minutes):02d}")
    
        # Calculate the localtime difference between consecutive rows
        data['localtime_Diff'] = data['DATETIME'].diff().dt.total_seconds().fillna(0)

    ###############
        data['Speed_kmh'] = data['RPM'] * 0.0836
        
        # Convert Speed to m/s
        data['Speed_ms'] = data['Speed_kmh'] / 3.6
        
        # Initialize total distance covered
        distance_per_mode = defaultdict(float)

        total_distance_with_RPM = 0

        for index, row in data.iterrows():
            if row['RPM'] > 0:

                distance_interval = row['Speed_ms'] * row['localtime_Diff']
                # Calculate the distance traveled in this interval
                total_distance_with_RPM += distance_interval
                
        print("Distance With RPM---------------------->",total_distance_with_RPM/1000)
    #################

    ###########Calculating the Distance based on Ground Distance from GPS Module data
        # total_distance_Ground_Distance = data['GROUND_DISTANCE'].iloc[-1]
        # Drop any empty values in the 'GROUND_DISTANCE' column and get the last non-empty value
        total_distance_Ground_Distance = data['GROUND_DISTANCE'].dropna().iloc[-1]

        # for index, row in data.iterrows():
        #     if row['GROUND_DISTANCE'] >= 100:

        #         distance_interval_groundSpeed = row['GROUND_DISTANCE'] * row['localtime_Diff']
        #         # Calculate the distance traveled in this interval
        #         total_distance_Ground_Distance += distance_interval_groundSpeed
                
        print("Distance With GROUND_DISTANCE---------------------->",total_distance_Ground_Distance/1000)
    ###############
    
    
        #Set the 'localtime' column as the index
        data.set_index('DATETIME', inplace=True)
    
        # duplicates = data.index[data.index.duplicated()]
        # print("Duplicate Index Values:", duplicates)
    
        # data = data.groupby(data.index).mean()  # Example: Taking the mean of duplicate values at the same timestamp
    
        # Keep the first value in each group of duplicates
        data = data.groupby(data.index).first()
    
        #Resample the data to have one-second intervals and fill missing values with previous ones
        data_resampled = data.resample('s').ffill()
    
    
        # Calculate the localtime difference between consecutive rows
        data_resampled['localtime_Diff'] = data_resampled.index.to_series().diff().dt.total_seconds().fillna(0)
    
        # Calculate the actual Ampere-hours (Ah) using the trapezoidal rule for numerical integration
    
    
        # print(data_resampled['Current_value'])
        actual_ah = abs((data_resampled['Current_value'] * data_resampled['localtime_Diff']).sum()) / 3600  # Convert seconds to hours
        print("Actual Ampere-hours (Ah):------------------------------------> {:.2f}".format(actual_ah))
    
        # Calculate the actual Watt-hours (Wh) using the trapezoidal rule for numerical integration
        watt_h = abs((data_resampled['Current_value'] * data_resampled['voltage_value'] * data_resampled['localtime_Diff']).sum()) / 3600  # Convert seconds to hours
        print("Actual Watt-hours (Wh):------------------------------------> {:.2f}" .format(watt_h))
    
    ####Commented because soc ah is not available in LXS dbc
        #starting and ending ah
        # starting_soc_Ah = data['SOCAh [SA: 08]'].iloc[0]
        # ending_soc_Ah = data['SOCAh [SA: 08]'].iloc[-1]
    
        # print("Starting SoC (Ah):{:.2f}".format (starting_soc_Ah))
        # print("Ending SoC (Ah):{:.2f}".format  (ending_soc_Ah))

###################

         # filtered_data_DischargeCurrent= data_resampled[-200 < data_resampled['Current_value'] < 0]                   #Only Discharge
        filtered_data_DischargeCurrent = data_resampled[(data_resampled['Regeneration']==0)]
        Discharge_ah = abs((filtered_data_DischargeCurrent['Current_value'] * filtered_data_DischargeCurrent['localtime_Diff']).sum()) / 3600  # Convert seconds to hours

        # Calculate the actual Watt-hours (Wh) using the trapezoidal rule for numerical integration
        watt_h = abs((data_resampled['Current_value'] * data_resampled['voltage_value'] * data_resampled['localtime_Diff']).sum()) / 3600  # Convert seconds to hours
        print("Actual Watt-hours (Wh):------------------------------------> {:.2f}" .format(watt_h))

        filtered_data_DischargeCurrent['localtime_Diff'] = data_resampled.index.to_series().diff().dt.total_seconds().fillna(0)
        discharge_wh_hr =abs((filtered_data_DischargeCurrent['Current_value'] * filtered_data_DischargeCurrent['voltage_value'] * filtered_data_DischargeCurrent['localtime_Diff']).sum()) / 3600  # Convert seconds to hours
        print("discharge Watt-hours (Wh):------------------------------------> {:.2f}" .format(discharge_wh_hr))


        #####calculating the distance with discharge
        filtered_data_DischargeCurrent['Speed_kmh'] = filtered_data_DischargeCurrent['RPM'] * 0.0836
        
        # Convert Speed to m/s
        filtered_data_DischargeCurrent['Speed_ms'] = filtered_data_DischargeCurrent['Speed_kmh'] / 3.6
        

        total_distance_with_RPM_discharge = 0

        for index, row in filtered_data_DischargeCurrent.iterrows():
            # if row['RPM'] > 0:
            if 0 < row['RPM'] < 1000:
                distance_interval = row['Speed_ms'] * row['localtime_Diff']
                # Calculate the distance traveled in this interval
                total_distance_with_RPM_discharge += distance_interval
                
        total_distance_with_RPM_discharge = total_distance_with_RPM_discharge/1000
        print("total_distance_with_RPM_discharge---------------------------------------------------------------------------------->",total_distance_with_RPM_discharge)

        
        # avg_discharge_current = filtered_data_DischargeCurrent['Current_value'].mean()
            # Filter out zeros
        non_zero_discharge_current = filtered_data_DischargeCurrent[filtered_data_DischargeCurrent['Current_value'] != 0]

        # Calculate the average of non-zero values
        avg_discharge_current = non_zero_discharge_current['Current_value'].mean()
        print("Average Discharge Current (ignoring zeros):", avg_discharge_current)


        filtered_data_RegenCurrent = data_resampled[(data_resampled['Regeneration']==1)]
        non_zero_regen_current = filtered_data_RegenCurrent[filtered_data_RegenCurrent['Current_value'] != 0]

        avg_regen_current = non_zero_regen_current['Current_value'].mean()
###################
    
    
        #Code for SOC_percentage(Starting and Ending SOC)
        starting_soc_percentage = data['SOC_value'].max()
        ending_soc_percentage = data['SOC_value'].min()
        print("Starting SOC:", starting_soc_percentage)
        print("Ending SOC:", ending_soc_percentage)
    
    
        # Initialize total distance covered
        total_distance = 0
        distance_per_mode = defaultdict(float)
    
        # Iterate over rows to compute distance covered between consecutive points
        for i in range(len(data) - 1):
        
            # Get latitude and longitude of consecutive points
            lat1 = data['LATITUDE'].iloc[i]
            lon1 = data['LONGITUDE'].iloc[i]
            lat2 = data['LATITUDE'].iloc[i + 1]
            lon2 = data['LONGITUDE'].iloc[i + 1]
        
    
            # Calculate distance between consecutive points
            distance = haversine(lat1, lon1, lat2, lon2)
    
            # Add distance to total distance covered
            total_distance += distance
    
            # mode = data['Mode_Ack [SA: 02]'].iloc[i]
            # distance_per_mode[mode] += distance
    
        print("Total distance covered (in kilometers):{:.2f}".format(total_distance))
    
        ##############   Wh/Km
        Wh_km = abs(watt_h / total_distance)
        print("WH/KM:{:.2f}". format (watt_h / total_distance))
    
        # Assuming 'data' is your DataFrame with 'SOC_value' column
        initial_soc = data['SOC_value'].iloc[-1]  # Initial SOC percentage
        final_soc = data['SOC_value'].iloc[0]   # Final SOC percentage
    
        # Calculate total SOC consumed
        total_soc_consumed =  abs(final_soc - initial_soc)
    
        print ("Total SOC consumed:{:.2f}".format (total_soc_consumed),"%")
    
    
        # # Check if the mode remains constant or changes
        # mode_values = data['Mode_Ack [SA: 02]'].unique()
    
    
    
    
        # Count number of 'ones' in 'eco' column
        eco_count = data_resampled['Eco_mode'].sum()
    
        # Count number of 'ones' in 'power' column
        power_count = data_resampled['power_mode'].sum()
    
        # Calculate total number of rows
        total_rows = len(data_resampled)
    
        # Calculate percentages
        eco_percentage = (eco_count / total_rows) * 100
        power_percentage = (power_count / total_rows) * 100
    
        # Display or use the results as needed
        print(f"Number of 'ones' in 'eco': {eco_count}")
        print(f"Number of 'ones' in 'power': {power_count}")
        print(f"Percentage of 'ones' in 'eco': {eco_percentage:.2f}%")
        print(f"Percentage of 'ones' in 'power': {power_percentage:.2f}%")
    
    
    
        # Calculate Wh/km for the entire ride
        watt_h = abs((data_resampled['Current_value'] * data_resampled['voltage_value'] * data_resampled['localtime_Diff']).sum()) / 3600
        # wh_per_km_total = abs(watt_h / total_distance)    
    
    
    
        # Calculate power using Current_value and voltage_value
        data_resampled['Power'] = data_resampled['Current_value'] * data_resampled['voltage_value']
    
        # Find the peak power
        peak_power = data_resampled['Power'].max()
        print("Peak Power:", peak_power)
    
        average_current =data_resampled['Current_value'].mean()
    
        # Calculate the average power
        average_power = abs(data_resampled['Power'].mean())
        print("Average Power:", average_power)
        # ALTITUDE = None
        # # Calculate ALTITUDE if 'ALTITUDE' column exists in data_resampled
        # if 'ALTITUDE' in data_resampled.columns:
        #     ALTITUDE = data_resampled['ALTITUDE'].max()  # Assuming you want to find the maximum ALTITUDE
        #     print("Maximum ALTITUDE:", ALTITUDE)
        # else:
        #     print("ALTITUDE data not available.")
    
    
        # Calculate energy regenerated (in watt-hours)
        energy_regenerated = ((data_resampled[data_resampled['Power'] > 0]['Power']*data_resampled['localtime_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
    
        # Calculate total energy consumed (in watt-hours)
        total_energy_consumed =  ((data_resampled[data_resampled['Power'] < 0]['Power']*data_resampled['localtime_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
    
        print("total",total_energy_consumed)
        print("total energy regenerated",energy_regenerated)
    
        # Calculate regenerative effectiveness as a percentage
        if total_energy_consumed != 0:
            print("-----------------------------------------------------------")
            regenerative_effectiveness = abs(energy_regenerated / total_energy_consumed) * 100
            print("Regenerative Effectiveness (%):", regenerative_effectiveness)
        else:
            regenerative_effectiveness = 0
            print("Total energy consumed is 0, cannot calculate regenerative effectiveness.")
    
        # Calculate idling localtime percentage (RPM was zero for more than 5 seconds)
        idling_localtime = ((data['RPM'] <= 0) | data['RPM'].isna()).sum()
        print("data length localtime",len(data))
        idling_percentage = (idling_localtime / len(data)) * 100
        print("Idling time percentage:", idling_percentage)
    
        # Calculate Time_specific speed ranges
        speed_ranges = [(0, 10), (10, 20), (20, 30), (30, 40), (40, 50),(50, 60),(60,70),(70, 80),(80, 90)]
        speed_range_percentages = {}
    
        for range_ in speed_ranges:
            speed_range_localtime = ((data['RPM'] * 0.0836 > range_[0]) & (data['RPM'] * 0.0836 < range_[1])).sum()
            speed_range_percentage = (speed_range_localtime / len(data)) * 100
            speed_range_percentages[f"Time_{range_[0]}-{range_[1]} km/h(%)"] = speed_range_percentage
            print(f"Time_{range_[0]}-{range_[1]} km/h(%): {speed_range_percentage:.2f}%")
    
            
        # Calculate power using Current_value and voltage_value
        data_resampled['Power'] = -data_resampled['Current_value'] * data_resampled['voltage_value']
    
        # Get the maximum cell voltage
        # max_cell_voltage = data_resampled['MaxCellVol [SA: 05]'].max()
    
        # Find the index where the maximum voltage occurs
        # max_index = data_resampled['MaxCellVol [SA: 05]'].idxmax()
    
        # Retrieve the corresponding cell ID using the index
        # max_cell_id = data_resampled['MaxVoltId [SA: 05]'].loc[max_index]
    
        # Get the minimum cell voltage
        # min_cell_voltage = data_resampled['MinCellVol [SA: 05]'].min()
    
        # Find the index where the minimum voltage occurs
        # min_index = data_resampled['MinCellVol [SA: 05]'].idxmin()
    
        # Retrieve the corresponding cell ID using the index
        # min_cell_id = data_resampled['MaxVoltId [SA: 05]'].loc[min_index]
    
        # voltage_difference = max_cell_voltage - min_cell_voltage
    
        # Get the maximum temperature
        # max_temp = data_resampled['MaxTemp [SA: 07]'].max()
    
        # # Find the index where the maximum temperature occurs
        # max_temp_index = data_resampled['MaxTemp [SA: 07]'].idxmax()
    
        # # Retrieve the corresponding temperature ID using the index
        # max_temp_id = data_resampled['MaxTempId [SA: 07]'].loc[max_temp_index]
    
        # Get the minimum temperature
        # min_temp = data_resampled['MinTemp [SA: 07]'].min()
    
        # # Find the index where the minimum temperature occurs
        # # min_temp_index = data_resampled['MinTemp [SA: 07]'].idxmin()
    
        # # # Retrieve the corresponding temperature ID using the index
        # # min_temp_id = data_resampled['MinTempId [SA: 07]'].loc[min_temp_index]
    
        # # Calculate the difference in temperature
        # temp_difference = max_temp - min_temp
    
        # # Get the maximum temperature of FetTemp [SA: 08]
        # max_fet_temp = data_resampled['FetTemp [SA: 08]'].max()
    
        # # Get the maximum temperature of AfeTemp [SA: 0C]
        # max_afe_temp = data_resampled['AfeTemp [SA: 0C]'].max()
    
        # # Get the maximum temperature of PcbTemp [SA: 0C]
        # max_pcb_temp = data_resampled['PcbTemp [SA: 0C]'].max()
    
        # # Get the maximum temperature of MCS_Temp
        # max_mcu_temp = data_resampled['MCS_Temp'].max()
    
        # # Check for abnormal motor temperature at high RPMs
        # max_motor_temp = data_resampled['Motor_Temp'].max()
    
        # Find the battery voltage
        # batteryVoltage = (data_resampled['voltage_value'].max()) * 10
        # print( "Battery Voltage", batteryVoltage )
    
        # Check for abnormal motor temperature at high RPMs for at least 15 seconds
        # abnormal_motor_temp = (data_resampled['Motor_Temp'] < 10) & (data_resampled['RPM'] > 3500)
        # abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
    
        # # Check if abnormal condition persists for at least 15 seconds
        # abnormal_motor_temp_detected = (abnormal_motor_temp_mask >= 120).any()
    
        # #For battery Analysis
        # cycleCount= data_resampled['CycleCount [SA: 07]'].max()
    
    
    
        # InfluxId= data['id'].iloc[0]
    
        max_continuous_duration = 0
        # data_resampled = data_resampled.dropna(subset=['RPM'])
        # data_resampled['RPM'] = data_resampled['RPM'].astype(int)
    
        # for speed in range(int(data_resampled['RPM'].min()), int(data_resampled['RPM'].max()) + 1):
        #     lower_bound = speed - window_size
        #     upper_bound = speed + window_size
        
    
        #     within_window = data_resampled[(data_resampled['RPM'] >= lower_bound) & (data_resampled['RPM'] <= upper_bound)].copy()
        #     within_window.loc[:, 'Group'] = (within_window['RPM'].diff() > window_size).cumsum()
        #     continuous_durations = within_window.groupby('Group')['localtime_Diff'].sum()
        
        #     current_max_duration = continuous_durations.max() if not continuous_durations.empty else 0
    
        #     if current_max_duration > max_continuous_duration:
        #         print("max_continuous_duration------->",max_continuous_duration)
        #         max_continuous_duration = current_max_duration
        #         cruising_rpm = speed
        #         cruising_speed=speed*0.0836
    
        #         if cruising_speed >1:
        #             cruise_speed=cruising_speed
        #         else :
        #             cruise_speed =0
    
    
            # Find the maximum value in the 'RPM' column
        Max_motor_rpm = data_resampled['RPM'].max()
    
        avg_motor_rpm =data_resampled['RPM'].mean()
    
        # Print the maximum motor speed in RPM
        print("The maximum motor speed in RPM is:", Max_motor_rpm)
    
        # Convert the maximum motor speed to speed using the given factor
        peak_speed = Max_motor_rpm * 0.0836
    
        avg_speed =avg_motor_rpm * 0.0836
    
        # Print the maximum speed
        print("The maximum speed is:", peak_speed)
    
        # total_energy_kwh = actual_ah * batteryVoltage / 1000
        # print("Total energy charged in kWh: {:.2f}".format(total_energy_kwh))
    
        # total_energy_kw = total_energy_kwh / total_duration.seconds / 3600
        # print("Electricity consumption units in kW", (total_energy_kw))
    
        # Add these variables and logic to ppt_data

        ending_soc_rows = data[data['SOC_value'] == ending_soc_percentage]


        if not ending_soc_rows.empty:
            ending_soc_first_occurrence = ending_soc_rows.iloc[0]
            ending_battery_voltage = ending_soc_first_occurrence['voltage_value']
            ending_battery_voltage = ending_battery_voltage
        else:
             # If exact starting SOC percentage is not found, find the nearest SOC percentage
            nearest_soc_index = (data['SOC_value'] - ending_soc_percentage).abs().idxmin()
            ending_soc_first_occurrence = data.iloc[nearest_soc_index]
            ending_battery_voltage = ending_soc_first_occurrence['voltage_value']
            ending_battery_voltage = ending_battery_voltage


        ###########For initial Temperature
        Starting_soc_rows = data[data['SOC_value'] == starting_soc_percentage]


        if not Starting_soc_rows.empty:
            Starting_soc_first_occurrence = Starting_soc_rows.iloc[0]
            Initial_MCU_TEMP = Starting_soc_first_occurrence['MCS_Temp']
            Initial_MOTOR_TEMP = Starting_soc_first_occurrence['Motor_Temp']
            Initial_BATTERY_TEMP = Starting_soc_first_occurrence['Battery_Temperature']
            

            
        else:
             # If exact starting SOC percentage is not found, find the nearest SOC percentage
            nearest_soc_index = (data['SOC_value'] - starting_soc_percentage).abs().idxmin()
            Initial_MCU_TEMP = Starting_soc_first_occurrence['MCS_Temp']
            Initial_MOTOR_TEMP = Starting_soc_first_occurrence['Motor_Temp']
            Initial_BATTERY_TEMP = Starting_soc_first_occurrence['Battery_Temperature']

        
        # Get the maximum temperature of MCS_Temp
        max_mcu_temp = data_resampled['MCS_Temp'].max()
        avg_mcu_temp = data_resampled['MCS_Temp'].mean()

            # Check for abnormal motor temperature at high RPMs
        max_motor_temp = data_resampled['Motor_Temp'].max()
        avg_motor_temp = data_resampled['Motor_Temp'].mean()

        max_BATTERY_temp = data_resampled['Battery_Temperature'].max()
        avg_BATTERY_temp = data_resampled['Battery_Temperature'].mean()

        
            
    
        ppt_data = {
            "Date and localtime": str(start_localtime) + " to " + str(end_localtime),
            # "INFLUX ID ": InfluxId,
            "Total Time taken for the ride": total_duration,
            # "Starting SoC (Ah)": starting_soc_Ah,
            # "Ending SoC (Ah)": ending_soc_Ah,
            "Actual Ampere-hours (Ah)": actual_ah,
            "Starting SoC (%)": starting_soc_percentage,
            "Ending SoC (%)": ending_soc_percentage,
            "Total SOC consumed(%)":starting_soc_percentage- ending_soc_percentage,
            "Energy consumption Rate(Wh/Km)": watt_h / total_distance,
            "Discharge Efficiency (Wh/km)":discharge_wh_hr/total_distance_with_RPM_discharge,
            "Total distance - RPM": total_distance_with_RPM/1000,
            "Total distance covered (km) - Lat & Long(GPS)": total_distance,
            "Total distance - Ground Distance(GPS)": total_distance_Ground_Distance/1000,
            # "Mode": "",
            # "Wh/km in Normal mode": wh_per_km_Normal_mode,
            # "Distance travelled in Normal mode":distance_per_mode[2],
            # "Wh/km in Fast Mode": wh_per_km_Fast_mode,
            # "Distance travelled in Fast Mode":distance_per_mode[6],
            # "Wh/km in ECO Mode": wh_per_km_ECO_mode,
            # "Distance travelled in ECO mode":distance_per_mode[4],
            # "Wh/km in LIMP Mode": wh_per_km_LIMP_mode,
            # "Distance travelled in LIMP Mode":distance_per_mode[5],
            "Eco_Mode_Percentage":eco_percentage,
            "Power_Mode_Percentage":power_percentage,
            "Actual Watt-hours (Wh)- Calculated_UsingFormala 'watt_h= 1/3600(|∑(V(t)⋅I(t)⋅Δt)|)'": watt_h,
            "Peak Power(W)": peak_power,
            "Average Power(W)": average_power,
            "Average_current":abs(average_current),
            "Average discharge current (A)":abs(avg_discharge_current),
            "Average Regen current (A)":abs(avg_regen_current),
            "Total Energy Regenerated(Wh)": energy_regenerated,
            # "Regenerative Effectiveness(%)": regenerative_effectiveness,
            "Avg_speed (km/hr)":avg_speed,
            "Voltage at cutoff (V)":ending_battery_voltage,
            "Initial Battery Temperature(C)": Initial_BATTERY_TEMP,
            "Maximum Battery Temperature(C)": max_BATTERY_temp,
            "Difference in Temperature(C)": max_BATTERY_temp- Initial_BATTERY_TEMP,
            "Initial MCU Temperature (at 100 SOC):":Initial_MCU_TEMP,
            "Maximum MCU Temperature(C)": max_mcu_temp,
            "Initial Motor Temperature (at 100 SOC)":Initial_MOTOR_TEMP,
            "Maximum Motor Temperature(C)": max_motor_temp,
            # "Highest Cell Voltage(V)": max_cell_voltage,
            # "Lowest Cell Voltage(V)": min_cell_voltage,
            # "Difference in Cell Voltage(V)": voltage_difference,
            # "Minimum Temperature(C)": min_temp,
            # "Maximum Temperature(C)": max_temp,
            # "Difference in Temperature(C)": max_temp- min_temp,
            # "Maximum Fet Temperature-BMS(C)": max_fet_temp,
            # "Maximum Afe Temperature-BMS(C)": max_afe_temp,
            # "Maximum PCB Temperature-BMS(C)": max_pcb_temp,
            # "Maximum MCU Temperature(C)": max_mcu_temp,
            # "Maximum Motor Temperature(C)": max_motor_temp,
            # "Abnormal Motor Temperature Detected(C)": abnormal_motor_temp_detected,
            # "highest cell temp(C)": max_cell_temp,
            # "lowest cell temp(C)": min_cell_temp,
            # "Difference between Highest and Lowest Cell Temperature at 100% SOC(C)": CellTempDiff,
            # "Battery Voltage(V)": batteryVoltage,
            # "Total energy charged(kWh)- Calculated_BatteryData": total_energy_kwh,
            # "Electricity consumption units(kW)": total_energy_kw,
            # "Cycle Count of battery": cycleCount,
            # "Cruising Speed (Rpm)": cruising_rpm,
            # "cruising_speed (km/hr)":cruise_speed,
            "Maximum Motor speed (RPM)":Max_motor_rpm,
            "Peak speed (Km/hr)":peak_speed,
            }
    
        # mode_values = data_resampled['Mode_Ack [SA: 02]'].unique() #If Mode_Ack [SA: 02] has values [1, 2, 2, 3, 1], unique() will return array([1, 2, 3]).
    
    #         mode  if len(mode_values) == 1:   # Mode remains constant throughout the log file
    #    = mode_values[0]
    #         if mode == 2:
    #             ppt_data["Mode"] = "Normal mode: 100%"
    #         elif mode == 6:
    #             ppt_data["Mode"] = "Fast Mode: 100%"
    #         elif mode == 4:
    #             ppt_data["Mode"] = "Eco mode: 100%"
    #         elif mode == 5:
    #             ppt_data["Mode"] = "Limp mode: 100%"
    #     else:
    #         # Mode changes throughout the log file
    #         # mode_counts = data_resampled['Mode_Ack [SA: 02]'].value_counts(normalize=True) * 100
    #         total_rows = len(data)
    #         mode_counts = data['Mode_Ack [SA: 02]'].value_counts()
    #         mode_strings = []  # Initialize list to store mode strings
        
    #          # Exclude mode 0 or empty from total_rows calculation
    #         total_rows -= mode_counts.get(0, 0)  # Subtract rows where mode is 0 (if any)
    #         total_rows -= mode_counts.get('', 0)  # Subtract rows where mode is empty string (if any)
    
    #         for mode,count in mode_counts.items():
    #             if mode not in [0, '']:  # Exclude mode 0 and empty string from percentage calculation
    #                 percentage = (count / total_rows) * 100
    #                 if mode == 2:
    #                     mode_strings.append(f"Normal mode\n{percentage:.2f}%")
    #                 elif mode == 6:
    #                     mode_strings.append(f"Fast Mode\n{percentage:.2f}%")
    #                 elif mode == 4:
    #                     mode_strings.append(f"Eco mode\n{percentage:.2f}%")
    #                 elif mode == 5:
    #                    mode_strings.append(f"Limp mode\n{percentage:.2f}%")
    #         ppt_data["Mode"] = "\n".join(mode_strings)
    
    
    
        # Add calculated parameters to ppt_data
        ppt_data["Idling time percentage"] = idling_percentage
        ppt_data.update(speed_range_percentages)
    
        # Calculate power using Current_value and voltage_value
        data_resampled['Power'] = -data_resampled['Current_value'] * data_resampled['voltage_value']
    
    
    
        # Filter the DataFrame to include only positive current values
        positive_current = data_resampled[data_resampled['Current_value'] < 0]
    
        # Calculate the mean of the positive current values
        average_current = positive_current['Current_value'].mean()
    
        # average_current =data_resampled['Current_value'].mean()
    
    
    
    
        # # Get the maximum cell voltage
        # max_cell_voltage = data_resampled['MaxCellVol [SA: 05]'].max()
    
        # # Find the index where the maximum voltage occurs
        # max_index = data_resampled['MaxCellVol [SA: 05]'].idxmax()
    
        # # Retrieve the corresponding cell ID using the index
        # # max_cell_id = data_resampled['MaxVoltId [SA: 05]'].loc[max_index]
    
        #     # Get the minimum cell voltage
        # min_cell_voltage = data_resampled['MinCellVol [SA: 05]'].min()
    
        # # Find the index where the minimum voltage occurs
        # min_index = data_resampled['MinCellVol [SA: 05]'].idxmin()
    
        # # Retrieve the corresponding cell ID using the index
        # # min_cell_id = data_resampled['MaxVoltId [SA: 05]'].loc[min_index]
    
        # voltage_difference = max_cell_voltage - min_cell_voltage
    
    
        # # print("Lowest Cell Voltage:", min_cell_voltage, "V, Cell ID:", min_cell_id)
        # # print("Highest Cell Voltage:", max_cell_voltage, "V, Cell ID:", max_cell_id)
        # print("Difference in Cell Voltage:", voltage_difference, "V")
    
        # # Get the maximum temperature
        # max_temp = data_resampled['MaxTemp [SA: 07]'].max()
    
        # # Find the index where the maximum temperature occurs
        # max_temp_index = data_resampled['MaxTemp [SA: 07]'].idxmax()
    
        # # Retrieve the corresponding temperature ID using the index
        # # max_temp_id = data_resampled['MaxTempId [SA: 07]'].loc[max_temp_index]
    
    
        # # Get the minimum temperature
        # min_temp = data_resampled['MinTemp [SA: 07]'].min()
    
        # # Find the index where the minimum temperature occurs
        # min_temp_index = data_resampled['MinTemp [SA: 07]'].idxmin()
    
        # # Retrieve the corresponding temperature ID using the index
        # # min_temp_id = data_resampled['MinTempId [SA: 07]'].loc[min_temp_index]
        # # Calculate the difference in temperature
        # temp_difference = max_temp - min_temp
        # print("temp_difference------------------------>",temp_difference)
        # print("Total distance------------------------->",total_distance)
    
        # # # Print the information
        # # print("Maximum Temperature:", max_temp, "C, Temperature ID:", max_temp_id)
        # # print("Minimum Temperature:", min_temp, "C, Temperature ID:", min_temp_id)
        # print("Difference in Temperature:", temp_difference, "C")
    
        # # Get the maximum temperature of FetTemp [SA: 08]
        # max_fet_temp = data_resampled['FetTemp [SA: 08]'].max()
        # print("Maximum Fet Temperature:", max_fet_temp, "C")
    
        # # Get the maximum temperature of AfeTemp [SA: 0C]
        # max_afe_temp = data_resampled['AfeTemp [SA: 0C]'].max()
        # print("Maximum Afe Temperature:", max_afe_temp, "C")
    
        # # Get the maximum temperature of PcbTemp [SA: 0C]
        # max_pcb_temp = data_resampled['PcbTemp [SA: 0C]'].max()
        # print("Maximum PCB Temperature:", max_pcb_temp, "C")
    
        # # Get the maximum temperature of MCS_Temp
        # max_mcu_temp = data_resampled['MCS_Temp'].max()
        # print("Maximum MCU Temperature:", max_mcu_temp, "C")
    
        # # Check for abnormal motor temperature at high RPMs
        # max_motor_temp = data_resampled['Motor_Temp'].max()
    
    
        # print("Maximum Motor Temperature:", max_motor_temp, "C")
    
    
        # # Check for abnormal motor temperature at high RPMs for at least 15 seconds
        # abnormal_motor_temp = (data_resampled['Motor_Temp'] < 10) & (data_resampled['RPM'] > 3500)
    
        # # Convert to a binary mask indicating consecutive occurrences
        # abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
    
        return total_duration, total_distance, Wh_km, total_soc_consumed,ppt_data
    
    
    
    # def capture_analysis_output(log_file, km_file, folder_path):
    def capture_analysis_output(log_file,folder_path):
        try:
            # Capture print statements
            analysis_output = io.StringIO()
            output_file = "analysis_results.docx"
    
            with redirect_stdout(analysis_output):
                #total_duration, total_distance, Wh_km, total_soc_consumed, ppt_data = analysis_Energy(log_file, km_file)
                analysis_output = analysis_output.getvalue()
            # Extract folder name from folder_path
            folder_name = os.path.basename(folder_path)
    
            # Create a new PowerPoint presentation
            prs = Presentation()
            # Add title slide with 'Selawik' style
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = f"Analysis Results from Folder - {folder_name}"
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.size = Pt(36)  # Corrected to Pt
            title.text_frame.paragraphs[0].font.name = 'Selawik'
    
            rows = len(ppt_data) + 1
            cols = 2
            table_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(table_slide_layout)
            shapes = slide.shapes
            # title_shape = shapes.title
            # title_shape.text = "Results:"
    
            # # Centering the title horizontally
            # title_shape.left = Inches(0.5)
            # title_shape.top = Inches(0.5)  # Adjust as needed
    
            # # Setting the font size of the title
            # title_shape.text_frame.paragraphs
    
            # Define maximum number of rows per slide
            max_rows_per_slide = 13
            # Add some space between title and table
            # title_shape.top = Inches(0.5)
    
            table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
            table.columns[0].width = Inches(4)
            table.columns[1].width = Inches(4)
            table.cell(0, 0).text = "Metric"
            table.cell(0, 1).text = "Value"
    
            # Initialize row index
            row_index = 1
            # Iterate over data and populate the table
            for key, value in ppt_data.items():
                # Check if current slide has reached maximum rows
                if row_index > max_rows_per_slide:
                    # Add a new slide
                    slide = prs.slides.add_slide(table_slide_layout)
                    shapes = slide.shapes
                    title_shape = shapes.title
                    title_shape.text = "Analysis Results:"
    
                    # Add a new table to the new slide
                    table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
                    table.columns[0].width = Inches(4)
                    table.columns[1].width = Inches(4)
                    table.cell(0, 0).text = "Metric"
                    table.cell(0, 1).text = "Value"
    
                    # Reset row index for the new slide
                    row_index = 1
    
                # Populate the table
                table.cell(row_index, 0).text = key
                table.cell(row_index, 1).text = str(value)
    
                # Increment row index
                row_index += 1
    
            # plot_ghps(data,folder_path)
    
            # Add image slide with title and properly scaled image
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
    
            # Remove the unwanted title placeholder
            for shape in slide.shapes:
                if shape.is_placeholder:
                    slide.shapes._spTree.remove(shape._element)
            # print("After function")
            # # Add the title
            # title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), prs.slide_width - Inches(2), Inches(1))
            # title_shape.text = "Graph Analysis"
            # print("After function 2")
            # Add the image and adjust its position and size
            # graph_width = prs.slide_width - Inches(1)
            # graph_height = prs.slide_height - Inches(2)
            # left = (prs.slide_width - graph_width) / 2
            # top = (prs.slide_height - graph_height) / 2 + Inches(1)
            # pic = slide.shapes.add_picture(graph_path, left, top, width=graph_width, height=graph_height)
            # Save the presentation
            output_file_name = f"{folder_path}/analysis_{folder_name}.pptx"
            prs.save(output_file_name)
            print("PPT generated!")
    
            # Create a new Excel workbook
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = "Analysis Results"
    
            # Populate Excel sheet with analysis data
            for i, (key, value) in enumerate(ppt_data.items(), start=1):
                ws.cell(row=i, column=1, value=key)
                ws.cell(row=i, column=2, value=value)
    
            # Save the Excel workbook
            excel_output_file = f"{folder_path}/analysis_{folder_name}.xlsx"
            wb.save(excel_output_file)
            print("excel_output_file1",excel_output_file," generated")
        
    
            df = pd.read_excel(excel_output_file, sheet_name="Analysis Results")
    
            # Transpose the DataFrame
            transposed_df = df.T
            transposed_df.columns = transposed_df.iloc[0]
            transposed_df = transposed_df[1:]
        
    
                # Define columns to plot for idling and speed metrics
            idling_speed_columns = [
                'Idling time percentage',
                'Time_0-10 km/h(%)',
                'Time_10-20 km/h(%)',
                'Time_20-30 km/h(%)',
                'Time_30-40 km/h(%)',
                'Time_40-50 km/h(%)',
                'Time_50-60 km/h(%)',
                'Time_60-70 km/h(%)',
                'Time_70-80 km/h(%)',
                # 'Time_80-90 km/h'
            ]
    
    
            # Define colors for idling and speed metrics
            idling_speed_colors = ['green', 'green', 'green', 'green', 'green', 'green', 'red', 'red', 'red']
    
    
            # Function to plot and save idling and speed metrics
            def plot_idling_speed_metrics():
                plt.figure(figsize=(15, 6))  # Increase figure size
                bars = plt.bar(idling_speed_columns, transposed_df.iloc[0][idling_speed_columns], color=idling_speed_colors)
                # plt.bar(idling_speed_columns, transposed_df.iloc[0][idling_speed_columns], color=idling_speed_colors)
                plt.xlabel('Metrics')
                plt.ylabel('Values')
                plt.title('Time_each speed interval')
                plt.xticks(rotation=-20, ha='left', fontsize=10)  # Adjust rotation and alignment
    
    
                # Annotate bars with values
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width() / 2, height, f'{height:.2f}', ha='center', va='bottom')
    
                plot_file = f"{folder_path}/Time spent in Speed intervals.png"
                plt.savefig(plot_file)
                # plt.show()
                print(f"Idling and speed bar plot saved: {plot_file}")
    
                # Insert plots into the Excel worksheet
                img_idling_speed = Image(plot_file)
                ws.add_image(img_idling_speed, 'F35')  # Adjust the cell location as needed
    
        
    
            # Generate and save the plots
            plot_idling_speed_metrics()
    
            # Save the Excel workbook
            excel_output_file = f"{folder_path}/analysis_{folder_name}.xlsx"
            wb.save(excel_output_file)
            print("excel_output_file1",excel_output_file," generated")
    
        
    
        except Exception as e:
            print("Error:", e)
    
    
    def current_percentage_calc(data,save_path):
        # Define the bins and labels for categorization
        bins = [0, 10, 20, 30, 40, 50, 60, 70, 80, 90, 100, float('inf')]
        labels = ['0 to 10', '10 to 20', '20 to 30', '30 to 40', '40 to 50',
                '50 to 60', '60 to 70', '70 to 80', '80 to 90', '90 to 100', '> 100']
      

        # Cut the data into categories using the bins
        categories = pd.cut(data['Current_value'], bins=bins, labels=labels, right=False)

        # Add the category column to the data
        data['Category'] = categories

        # Calculate the sum of 'Current_value' values in each category
        category_sums = data.groupby('Category')['Current_value'].sum()

        # Calculate the absolute sum for percentage calculation
        total_sum = category_sums.abs().sum()

        # Calculate the percentage each category contributes
        percentages = (category_sums.abs() / total_sum * 100).reset_index(name='Percentage')

        # Plotting the results
        plt.figure(figsize=(14, 8))
        bars = plt.bar(percentages['Category'], percentages['Percentage'], color='skyblue')

        # Adding the percentage labels on top of each bar
        for bar in bars:
            yval = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2, yval + 0.5, f"{round(yval, 2)}%", ha='center', va='bottom')

        plt.xlabel('Current Value Ranges (A)')
        plt.ylabel('Percentage of Total Current (%)')
        plt.title('Percentage Distribution of Current Values')
        plt.xticks(rotation=90)  # Rotate labels for better visibility
        plt.grid(axis='y', linestyle='--', alpha=0.7)
        plt.tight_layout()
        # Save the plot to the specified path
        plot_file = f"{save_path}/Current Distribution.png"
        plt.savefig(plot_file)
        print("Current distribution Image saved")
        # plt.show()
    
    # Initialize variables to store file paths
    log_file = None
    
    main_folder_path = input_folder_path
    process_files_in_directory(main_folder_path) 
    
    def mergeExcel(main_folder_path):
        def prepare_sheet_in_memory(file_path):
            workbook = load_workbook(filename=file_path)
            sheet = workbook.active
            file_name = os.path.basename(file_path)
            sheet.insert_rows(1)
            sheet['A1'] = 'file name'
            sheet['B1'] = file_name
            return sheet, file_name
    
        def sheet_to_dict(sheet):
            data_dict = {}
            for row in sheet.iter_rows(min_row=2, values_only=True):
                key, *values = row
                data_dict[key] = values
            return data_dict
    
        def merge_dicts(dict1, dict2):
            for key, values in dict2.items():
                if key in dict1:
                    # Assuming that values is a list of values
                    dict1[key].extend(values)
                else:
                    dict1[key] = values
            return dict1
    
        def process_directory(directory):
            merged_data = {}
            file_names = []
            for root, dirs, files in os.walk(directory):
                for name in dirs:
                        subdir_path = os.path.join(root, name)
                        for file_name in os.listdir(subdir_path):
                            if file_name.endswith(".xlsx"):
                                file_path = os.path.join(subdir_path, file_name)
                                sheet, extracted_file_name = prepare_sheet_in_memory(file_path)
                                data_dict = sheet_to_dict(sheet)
                                merged_data = merge_dicts(merged_data, data_dict)
                                if extracted_file_name not in file_names:
                                    file_names.append(extracted_file_name)
            return merged_data, file_names
    
        def merge_data_and_save_to_excel(main_folder_path):
            directory = main_folder_path
            merged_data, file_names = process_directory(directory)
    
            merged_workbook = Workbook()
            merged_sheet = merged_workbook.active
    
            # # Use the extracted file names for headers
            # headers = ['File name'] + file_names
            # merged_sheet.append(headers)
    
            if merged_data:
                for key, values in merged_data.items():
                    merged_sheet.append([key] + values)
            else:
                print("No data found in merged_data")
    
            merged_file_path = os.path.join(directory, 'Analysis.xlsx')
            merged_workbook.save(filename=merged_file_path)
            print("Analysis file is ready")
    
        merge_data_and_save_to_excel(main_folder_path)
    
    
    
    
    # Iterate over immediate subfolders of main_folder_path
    for subfolder_1 in os.listdir(main_folder_path):
        subfolder_1_path = os.path.join(main_folder_path, subfolder_1)
    
        # Check if subfolder_1 is a directory
        if os.path.isdir(subfolder_1_path):
    
    # Iterate over subfolders within subfolder_1
            for subfolder in os.listdir(subfolder_1_path):
                subfolder_path = os.path.join(subfolder_1_path, subfolder)
                print(subfolder_path)
            
                # Check if subfolder starts with "Battery" and is a directory
                if os.path.isdir(subfolder_path):
                    log_file = None
                    log_found = False
                
                    # Iterate through files in the subfolder
                    for file in os.listdir(subfolder_path):
                        print(subfolder_path)
                        if file.startswith('log.') and file.endswith('.csv'):
                        
                            log_file = os.path.join(subfolder_path, file)
                            log_found = True
                            break  # Stop searching once the log file is found
                
                    # Process the log file if found
                    if log_found:
                        print(f"Processing log file: {log_file}")
                        try:
                            data = pd.read_csv(log_file)
                            # Process your data here
                        except Exception as e:
                            print(f"Error processing {log_file}: {e}")
            
            
                        total_duration = 0
                        total_distance = 0
                        Wh_km = 0
                        SOC_consumed = 0
                        mode_values = 0
                    
                        total_duration, total_distance, Wh_km, SOC_consumed, ppt_data = analysis_Energy(data,subfolder_path)
                        capture_analysis_output(log_file, subfolder_path)
                        current_percentage_calc(data,subfolder_path)
                    
                    else:
                        print("Log file file not found in subfolder:", subfolder)
    
    mergeExcel(main_folder_path)