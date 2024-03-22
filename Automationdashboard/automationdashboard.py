import pandas as pd
from math import radians, sin, cos, sqrt, atan2
 
import sys
import io
 
import matplotlib.pyplot as plt
from mpl_toolkits.basemap import Basemap
 
from PIL import Image
import os
import matplotlib.dates as mdates
import mplcursors  # Import mplcursors
 
from contextlib import redirect_stdout
from pptx.util import Inches, Pt  # Correcting the import statement
 
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches
# Path to the folder containing the CSV files
path = r"c:\Users\annmon.james\lectrix_internship\work\Automationdashboard\data"
 
# List to store DataFrames from each CSV file
dfs = []
def haversine(lat1, lon1, lat2, lon2):
    """
    Calculate the great-circle distance between two points
    on the Earth's surface using the Haversine formula.
    """
    # Convert latitude and longitude from degrees to radians
    lat1 = radians(lat1)
    lon1 = radians(lon1)
    lat2 = radians(lat2)
    lon2 = radians(lon2)
 
    # Radius of the Earth in kilometers
    R = 6371.0
   
    # Calculate differences in latitude and longitude
    dlat = lat2 - lat1
    dlon = lon2 - lon1
 
    # Calculate Haversine formula
    a = sin(dlat / 2)**2 + cos(lat1) * cos(lat2) * sin(dlon / 2)**2
    c = 2 * atan2(sqrt(a), sqrt(1 - a))
    distance = R * c
 
    return distance
# Define a function to set current to zero if RPM is zero for 10 or more consecutive points
def adjust_current(row):
    adjust_current.zero_count = getattr(adjust_current, 'zero_count', 0)
    if row['MotorSpeed_340920578'] == 0:
        adjust_current.zero_count += 1
    else:
        adjust_current.zero_count = 0
   
    if adjust_current.zero_count >= 10:
        return 0
    else:
        return row['PackCurr_6']
def plot_ghps(log_file):
 
    data = pd.read_csv(r"c:\Users\annmon.james\lectrix_internship\work\Automationdashboard\data\log_file.csv")
 
    # Apply the adjustment function to the DataFrame
 
 
    data['localtime'] = pd.to_datetime(data['localtime'])
    data.set_index('localtime', inplace=True)
    data['PackCurr_6'] = data.apply(adjust_current, axis=1)
 
    # Create a figure and axes for plotting
    fig, ax1 = plt.subplots(figsize=(10, 6))
 
    # Plot 'PackCurr_6' on primary y-axis
    line1, = ax1.plot(data.index, -data['PackCurr_6'], color='blue', label='PackCurr_6')
    ax1.set_ylabel('Pack Current (A)', color='blue')
    ax1.yaxis.set_label_coords(-0.1, 0.7)  # Adjust label position
 
    # Create secondary y-axis for 'MotorSpeed_340920578' (RPM)
    ax2 = ax1.twinx()
    line2, = ax2.plot(data.index, data['MotorSpeed_340920578'], color='green', label='Motor Speed')
    ax2.set_ylabel('Motor Speed (RPM)', color='green')
 
    # Add 'AC_Current_340920579' to primary y-axis
    line3, = ax1.plot(data.index, data['AC_Current_340920579'], color='red', label='AC Current')
 
    # Add 'AC_Voltage_340920580' scaled to 10x to the left side y-axis
    line4, = ax1.plot(data.index, data['AC_Voltage_340920580'] * 10, color='orange', label='AC Voltage (x10)')
 
    # Add 'Throttle_408094978' to the left side y-axis
    line5, = ax1.plot(data.index, data['Throttle_408094978'], color='lightgray', label='Throttle (%)')
 
    # Hide the y-axis label for 'AC_Current_340920579'
    ax1.get_yaxis().get_label().set_visible(False)
 
    # Set x-axis label and legend
    ax1.set_xlabel('Local Time')
    ax1.legend(loc='upper left')
    ax2.legend(loc='upper right')
 
    # Add a title to the plot
    plt.title('Battery Pack, Motor Data, and Throttle')
 
    # Format x-axis ticks as hours:minutes:seconds
    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M:%S'))
 
    # Set grid lines lighter
    ax1.grid(True, linestyle=':', linewidth=0.5, color='gray')
    ax2.grid(True, linestyle=':', linewidth=0.5, color='gray')
 
    # Enable cursor to display values on graphs
    mplcursors.cursor([line1, line2, line3, line4, line5])
 
    # Save the plot as an image or display it
    plt.tight_layout()  # Adjust layout to prevent clipping of labels
    plt.savefig('graph.png')  # Save the plot as an image
    plt.show()
def analysis_Energy(log_file, km_file):
 
    data = pd.read_csv(r"c:\Users\annmon.james\lectrix_internship\work\Automationdashboard\data\log_file.csv")
 
    data_KM = pd.read_csv(r"c:\Users\annmon.james\lectrix_internship\work\Automationdashboard\data\km_file.csv")
 
    total_duration = 0
    total_distance = 0
    Wh_km = 0
    SOC_consumed = 0
 
    # Check if 'localtime' column exists in data DataFrame
    if 'localtime' not in data.columns:
        print("Error: 'localtime' column not found in the DataFrame.")
        return None, None, None, None
 
 
 
 
    # Drop rows with missing values in 'SOCAh_8' column
    data.dropna(subset=['SOCAh_8'], inplace=True)
 
    # Convert the 'localtime' column to datetime format
    data['localtime'] = pd.to_datetime(data['localtime'], format='%d/%m/%Y %H:%M:%S.%f')
 
    # Calculate the start time and end time
    start_time = data['localtime'].min()
    end_time = data['localtime'].max()
 
    # Calculate the total time taken for the ride
    total_duration = end_time - start_time
 
    total_hours = total_duration.seconds // 3600
    total_minutes = (total_duration.seconds % 3600) // 60
 
    print(f"Total time taken for the ride: {int(total_hours):02d}:{int(total_minutes):02d}")
 
    # Calculate the time difference between consecutive rows
    data['Time_Diff'] = data['localtime'].diff().dt.total_seconds().fillna(0)
 
    # Set the 'localtime' column as the index
    data.set_index('localtime', inplace=True)
 
    # Resample the data to have one-second intervals and fill missing values with previous ones
    data_resampled = data.resample('S').ffill()
 
    # Calculate the time difference between consecutive rows
    data_resampled['Time_Diff'] = data_resampled.index.to_series().diff().dt.total_seconds().fillna(0)
 
    # Calculate the actual Ampere-hours (Ah) using the trapezoidal rule for numerical integration
    actual_ah = (data_resampled['PackCurr_6'] * data_resampled['Time_Diff']).sum() / 3600  # Convert seconds to hours
    print("Actual Ampere-hours (Ah):", actual_ah)
 
    # Calculate the actual Watt-hours (Wh) using the trapezoidal rule for numerical integration
    watt_h = (data_resampled['PackCurr_6'] * data_resampled['PackVol_6'] * data_resampled['Time_Diff']).sum() / 3600  # Convert seconds to hours
    print("Actual Watt-hours (Wh):", watt_h)
 
    ###########   starting and ending ah
    starting_soc = data['SOCAh_8'].iloc[0]
    ending_soc = data['SOCAh_8'].iloc[-1]
 
    print("Starting SoC (Ah):", starting_soc)
    print("Ending SoC (Ah):", ending_soc)
 
    ##################### KM ---------------------
    # Convert the 'localtime' column to datetime format
    data_KM['localtime'] = pd.to_datetime(data_KM['localtime'])
 
    # Initialize total distance covered
    total_distance = 0
 
    # Iterate over rows to compute distance covered between consecutive points
    for i in range(1, len(data_KM)):
        # Get latitude and longitude of consecutive points
        lat1, lon1 = data_KM.loc[i - 1, 'latitude'], data_KM.loc[i - 1, 'longitude']
        lat2, lon2 = data_KM.loc[i, 'latitude'], data_KM.loc[i, 'longitude']
 
        # Calculate distance between consecutive points
        distance = haversine(lat1, lon1, lat2, lon2)
 
        # Add distance to total distance covered
        total_distance += distance
 
    print("Total distance covered (in kilometers):", total_distance)
 
    ##############   Wh/Km
    Wh_km = watt_h / total_distance
    print("WH/KM       :", watt_h / total_distance)
 
        # Assuming 'data' is your DataFrame with 'SOC_8' column
    initial_soc = data['SOC_8'].iloc[0]  # Initial SOC percentage
    final_soc = data['SOC_8'].iloc[-1]   # Final SOC percentage
 
    # Calculate total SOC consumed
    total_soc_consumed =  final_soc - initial_soc
 
    print("Total SOC consumed:", total_soc_consumed, "%")
 
 
 
    # Check if the mode remains constant or changes
    mode_values = data['Mode_Ack_408094978'].unique()
 
    if len(mode_values) == 1:
        # Mode remains constant throughout the log file
        mode = mode_values[0]
        if mode == 3:
            print("Mode is Custom mode.")
        elif mode == 2:
            print("Mode is Sports mode.")
        elif mode == 1:
            print("Mode is Eco mode.")
    else:
        # Mode changes throughout the log file
        mode_counts = data['Mode_Ack_408094978'].value_counts(normalize=True) * 100
        for mode, percentage in mode_counts.items():
            if mode == 3:
                print(f"Custom mode: {percentage:.2f}%")
            elif mode == 2:
                print(f"Sports mode: {percentage:.2f}%")
            elif mode == 1:
                print(f"Eco mode: {percentage:.2f}%")
 
 
 
 
    return total_duration, total_distance, Wh_km, total_soc_consumed
folder_path = r"C:\Users\annmon.james\lectrix_internship\work\Automationdashboard\data"

def capture_analysis_output(log_file, km_file, folder_path):
    try:
        # Capture print statements
        analysis_output = io.StringIO()
        output_file = "analysis_results.docx"
 
        with redirect_stdout(analysis_output):
            total_duration, total_distance, Wh_km, total_soc_consumed = analysis_Energy(log_file, km_file)
        analysis_output = analysis_output.getvalue()
 
        # Extract folder name from folder_path
        folder_name = os.path.basename(folder_path)
 
        # Create a new PowerPoint presentation
        prs = Presentation()
 
        # Add title slide with 'Selawik' style
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"Analysis Results from Folder - {folder_name}"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)  # Corrected to Pt
        title.text_frame.paragraphs[0].font.name = 'Selawik'
 
        rows = 10
        cols = 2
        table_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(table_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Analysis Results:"
       
        # Add some space between title and table
        title_shape.top = Inches(0.5)
       
        table = shapes.add_table(rows, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(4)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        data = {
            "Total time taken for the ride": total_duration,
            "Actual Ampere-hours (Ah)": -0.496635,
            "Actual Watt-hours (Wh)": -25.81896295027778,
            "Starting SoC (Ah)": 33.91,
            "Ending SoC (Ah)": 34.433,
            "Total distance covered (in kilometers)": total_distance,
            "WH/KM": -40.360253303742816,
            "Total SOC consumed": "1.0%",
            "Mode": "Eco mode"
        }
        row_index = 1
        for key, value in data.items():
            table.cell(row_index, 0).text = key
            table.cell(row_index, 1).text = str(value)
            row_index += 1
 
 
        # Add image slide with title and properly scaled image
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
 
        # Remove the unwanted title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)
 
        # Add the title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), prs.slide_width - Inches(2), Inches(1))
        title_shape.text = "Graph Analysis"
 
        # Add the image and adjust its position and size
        graph_width = prs.slide_width - Inches(1)
        graph_height = prs.slide_height - Inches(2)
        left = (prs.slide_width - graph_width) / 2
        top = (prs.slide_height - graph_height) / 2 + Inches(1)
        pic = slide.shapes.add_picture('graph.png', left, top, width=graph_width, height=graph_height)
 
        # Save the presentation
        output_file_name = f"{folder_path}/analysis_{folder_name}.pptx"
        prs.save(output_file_name)
       
    except Exception as e:
        print("Error:", e)

#folder_path = "/home/sanjith/Documents/Graphs _ creta/15-51_16-00"
 
# Get the list of files in the folder
files = os.listdir(folder_path)
 
# Initialize variables to store file paths
log_file = None
km_file = None
 
 
 
 
# Path to the main folder containing subfolders
main_folder_path = r"c:\Users\annmon.james\lectrix_internship\work\Automationdashboard"
 
# Iterate over subfolders
for subfolder in os.listdir(main_folder_path):
    subfolder_path = os.path.join(main_folder_path, subfolder)
    print(subfolder)
    if os.path.isdir(subfolder_path):
        log_file = None
        km_file = None
        # Find 'log' and 'km' files
        l=0
        for file in os.listdir(subfolder_path):
            if file.startswith('log') and file.endswith('.csv'):
                log_file = os.path.join(subfolder_path, file)
                l = 1
            elif file.startswith('km') and file.endswith('.csv'):
                km_file = os.path.join(subfolder_path, file)
                l =2
 
 
# Read the CSV file into a pandas DataFrame
 
            if (l ==2):
                total_duration =0
                total_distance =0
                Wh_km =0
                SOC_consumed=0
 
 
                ### plot graphs
 
 
                plot_ghps(log_file)
                # total_duration, total_distance, Wh_km,SOC_consumed=analysis_Energy(log_file,km_file)
                capture_analysis_output(log_file, km_file, subfolder_path)