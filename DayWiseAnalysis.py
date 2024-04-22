
import pandas as pd
from math import radians, sin, cos, sqrt, atan2
import sys
import io
import openpyxl
import matplotlib.pyplot as plt
from mpl_toolkits.basemap import Basemap
from PIL import Image
import os
import matplotlib.dates as mdates
import mplcursors  # Import mplcursors
from contextlib import redirect_stdout
from pptx.util import Inches, Pt  # Correcting the import statement
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches
from openpyxl import load_workbook, Workbook
 
# Path to the folder containing the CSV files
# path = r"C:\Users\kamalesh.kb\CodeForAutomation\MAIN_FOLDER\MAR_21"
 
 
 
# List to store DataFrames from each CSV file
dfs = []
def haversine(lat1, lon1, lat2, lon2):
    """
    Calculate the great-circle distance between two points
    on the Earth's surface using the Haversine formula.
    """
    # Convert latitude and longitude from degrees to radians
    lat1 = radians(lat1)
    lon1 = radians(lon1)
    lat2 = radians(lat2)
    lon2 = radians(lon2)
 
    # Radius of the Earth in kilometers
    R = 6371.0
   
    # Calculate differences in latitude and longitude
    dlat = lat2 - lat1
    dlon = lon2 - lon1
 
    # Calculate Haversine formula
    a = sin(dlat / 2)**2 + cos(lat1) * cos(lat2) * sin(dlon / 2)**2
    c = 2 * atan2(sqrt(a), sqrt(1 - a))
    distance = R * c
 
    return distance
 
# Define a function to set current to zero if RPM is zero for 10 or more consecutive points
def adjust_current(row):
    adjust_current.zero_count = getattr(adjust_current, 'zero_count', 0)
    if row['MotorSpeed_340920578'] == 0:
        adjust_current.zero_count += 1
    else:
        adjust_current.zero_count = 0
   
    if adjust_current.zero_count >= 10:
        return 0
    else:
        return row['PackCurr_6']
   
def plot_ghps(data):


    if 'localtime' not in data.columns:
        # Convert the 'timestamp' column from Unix milliseconds to datetime
        data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')

        # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
        data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)


    # data = pd.read_csv(r"C:\Users\kamalesh.kb\CodeForAutomation\MAIN_FOLDER\MAR_21\log_file.csv")
 
    # Apply the adjustment function to the DataFrame
 
    # data['localtime'] = pd.to_datetime(data['localtime'], format='%d/%m/%Y %H:%M:%S.%f', dayfirst=True)
    data['localtime'] = pd.to_datetime(data['localtime'])
    data.set_index('localtime', inplace=True)
    data['PackCurr_6'] = data.apply(adjust_current, axis=1)
   
    # Create a figure and axes for plotting
    fig, ax1 = plt.subplots(figsize=(10, 6))
 
    # Plot 'PackCurr_6' on primary y-axis
    line1, = ax1.plot(data.index, -data['PackCurr_6'], color='blue', label='PackCurr_6')
    ax1.set_ylabel('Pack Current (A)', color='blue')
    ax1.yaxis.set_label_coords(-0.1, 0.7)  # Adjust label position
 
    # Create secondary y-axis for 'MotorSpeed_340920578' (RPM)
    ax2 = ax1.twinx()
    line2, = ax2.plot(data.index, data['MotorSpeed_340920578'], color='green', label='Motor Speed')
    ax2.set_ylabel('Motor Speed (RPM)', color='green')
 
    # Add 'AC_Current_340920579' to primary y-axis
    line3, = ax1.plot(data.index, data['AC_Current_340920579'], color='red', label='AC Current')
 
    # Add 'AC_Voltage_340920580' scaled to 10x to the left side y-axis
    line4, = ax1.plot(data.index, data['AC_Voltage_340920580'] * 10, color='yellow', label='AC Voltage (x10)')
 
    # Add 'Throttle_408094978' to the left side y-axis
    line5, = ax1.plot(data.index, data['Throttle_408094978'], color='orange', label='Throttle (%)')
 

    # Hide the y-axis label for 'AC_Current_340920579'
    ax1.get_yaxis().get_label().set_visible(False)
 
    # Set x-axis label and legend
    ax1.set_xlabel('Local Time')
    ax1.legend(loc='upper left')
    ax2.legend(loc='upper right')
 
    # Add a title to the plot
    plt.title('Battery Pack, Motor Data, and Throttle')
 
    # Format x-axis ticks as hours:minutes:seconds
    ax1.xaxis.set_major_formatter(mdates.DateFormatter('%H:%M:%S'))
 
    # Set grid lines lighter
    ax1.grid(True, linestyle=':', linewidth=0.5, color='gray')
    ax2.grid(True, linestyle=':', linewidth=0.5, color='gray')
 
    # Enable cursor to display values on graphs
    mplcursors.cursor([line1, line2, line3, line4, line5])
 
    # Save the plot as an image or display it
    plt.tight_layout()  # Adjust layout to prevent clipping of labels
    plt.savefig('graph.png')  # Save the plot as an image
    plt.show()
 
 
###################################
def splitBatteryWise(main_folder_path):
 
    ############
    for subfolder in os.listdir(main_folder_path):
        if subfolder.startswith("Battery"):  # Ignore folders starting with "Battery"
            continue
 
        subfolder_path = os.path.join(main_folder_path, subfolder)
        print(subfolder)
        if os.path.isdir(subfolder_path):
            log_file = None
            # km_file = None
            # Find 'log' and 'km' files
            log_found = False
            # km_found = False
            for file in os.listdir(subfolder_path):
                if file.startswith('log_withoutanamoly') and file.endswith('.csv'):
                    log_file = os.path.join(subfolder_path, file)
                    log_found = True
                # elif file.startswith('km') and file.endswith('.csv'):
                #     km_file = os.path.join(subfolder_path, file)
                #     km_found = True
                # if log_found and km_found:
                if log_found:
                    break
 
            # if log_found and km_found:
            if log_found:
                df= pd.read_csv(log_file)
               
 
# def analysis_Energy(log_file, km_file):
def analysis_Energy(log_file):
 
    print("Entered analysis energy")
    dayfirst=True
    data = pd.read_csv(log_file)

    if 'localtime' not in data.columns:
    # Convert the 'timestamp' column from Unix milliseconds to datetime
        data['localtime'] = pd.to_datetime(data['timestamp'], unit='ms')
 
    # Adjusting the timestamp to IST (Indian Standard Time) by adding 5 hours and 30 minutes
        data['localtime'] = data['localtime'] + pd.Timedelta(hours=5, minutes=30)
   
 
    total_duration = 0
    total_distance = 0
    Wh_km = 0
    SOC_consumed = 0
   
    if (data['SOC_8'] > 90).any():
        # Maximum cell temperature calculation
        temp_columns_max = [f'Temp{i}_10' for i in range(1, 9)]
        max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
        print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
 
        # Minimum cell temperature calculation
        temp_columns_min = [f'Temp{i}_10' for i in range(1, 9)]
        min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
        print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
 
        #Difference between Maximum and Minimum cell Temperature
        CellTempDiff= max_cell_temp-min_cell_temp
        print("Temperature difference: ",CellTempDiff)
 
    else:
        print("SOC is not maximum!")
         # Maximum cell temperature calculation
        temp_columns_max = [f'Temp{i}_10' for i in range(1, 9)]
        max_values = data[temp_columns_max].max(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        max_cell_temp = max_values.max()                  # Find the maximum among those maximum values
        print("\nOverall maximum value of cell temperature among those maximum values:", max_cell_temp)
 
        # Minimum cell temperature calculation
        temp_columns_min = [f'Temp{i}_10' for i in range(1, 9)]
        min_values = data[temp_columns_min].min(axis=1)     # Find the maximum value out of 8 columns (from Temp1_10 to Temp8_10)
        min_cell_temp = min_values.min()                  # Find the maximum among those maximum values
        print("\nOverall minimum value of cell temperature among those minimum values:", min_cell_temp)
 
        #Difference between Maximum and Minimum cell Temperature
        CellTempDiff= max_cell_temp-min_cell_temp
        print("Temperature difference: ",CellTempDiff)
 
 
   
 
 
    # Check if 'localtime' column exists in data DataFrame
    if 'localtime' not in data.columns:
        print("Error: 'localtime' column not found in the DataFrame.")
        return None, None, None, None
   
    # Drop rows with missing values in 'SOCAh_8' column
    data.dropna(subset=['SOCAh_8'], inplace=True)
 
    # Convert the 'localtime' column to datetime s
    # data['localtime'] = pd.to_datetime(data['localtime'], format='%d/%m/%Y %H:%M:%S.%f',dayfirst=True)
    data['localtime'] = pd.to_datetime(data['localtime'])
 
    # Calculate the start time and end time
    start_time = data['localtime'].min()
    end_time = data['localtime'].max()

    # Calculate the start time and end time with formatting
    start_time_seconds = data['localtime'].min().strftime('%d/%m/%Y %H:%M:%S')
    end_time_seconds = data['localtime'].max().strftime('%d/%m/%Y %H:%M:%S')
    
    print(start_time,end_time)

    # Calculate the total time taken for the ride``
    total_duration = end_time - start_time
    total_hours = total_duration.seconds // 3600
    total_minutes = (total_duration.seconds % 3600) // 60
 
    print(f"Total time taken for the ride: {int(total_hours):02d}:{int(total_minutes):02d}")
 
    # Calculate the time difference between consecutive rows
    data['Time_Diff'] = data['localtime'].diff().dt.total_seconds().fillna(0)
 
    # Set the 'localtime' column as the index
    data.set_index('localtime', inplace=True)
 
    # Resample the data to have one-second intervals and fill missing values with previous ones
    data_resampled = data.resample('s').ffill()
 
    # Calculate the time difference between consecutive rows
    data_resampled['Time_Diff'] = data_resampled.index.to_series().diff().dt.total_seconds().fillna(0)
 
    # Calculate the actual Ampere-hours (Ah) using the trapezoidal rule for numerical integration
    actual_ah = abs((data_resampled['PackCurr_6'] * data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours
    print("Actual Ampere-hours (Ah): {:.2f}".format(actual_ah))
 
    # Calculate the actual Watt-hours (Wh) using the trapezoidal rule for numerical integration
    watt_h = abs((data_resampled['PackCurr_6'] * data_resampled['PackVol_6'] * data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours
    print("Actual Watt-hours (Wh):{:.2f}" .format(watt_h))
 
    ###########   starting and ending ah
    starting_soc_Ah = data['SOCAh_8'].iloc[-1]
    ending_soc_Ah = data['SOCAh_8'].iloc[0]
 
    print("Starting SoC (Ah):{:.2f}".format (starting_soc_Ah))
    print("Ending SoC (Ah):{:.2f}".format  (ending_soc_Ah))
 
    # #Added date and time
    # starting_time= data['localtime'].iloc[-1]
   
    # Ending_time= data['localtime'].iloc[0]
    # print(starting_time)
    # print(Ending_time)


 ####################
 #the following code is providing wrong value for SOC_percentage, so skipped it.
    # Calculate starting and ending SoC in percentage
    # starting_soc_percentage = (starting_soc / actual_ah) * 100
    # ending_soc_percentage = (ending_soc / actual_ah) * 100
 
    # # Print the results
    # print("Starting SoC (%): {:.2f}%".format(starting_soc_percentage))
    # print("Ending SoC (%): {:.2f}%".format(ending_soc_percentage))
 #####################
 
 
 #Code for SOC_percentage(Starting and Ending SOC)
    starting_soc_percentage = data['SOC_8'].max()
    ending_soc_percentage = data['SOC_8'].min()
    print("Starting SOC:", starting_soc_percentage)
    print("Ending SOC:", ending_soc_percentage)
 
 
    ##################### KM ---------------------
    # Convert the 'localtime' column to datetime format
    # data_KM['localtime'] = pd.to_datetime(data_KM['localtime'], format='%d/%m/%Y %H:%M:%S.%f', dayfirst=True)
    # data['localtime'] = pd.to_datetime(data['localtime'])
   
    # Initialize total distance covered
    total_distance = 0
 
    # Iterate over rows to compute distance covered between consecutive points
    for i in range(len(data) - 1):
        # print("length--------------->",len(data))
        # print("i------------------->",i)
        # print(data['latitude'])
        # print(data['longitude'])
        # Get latitude and longitude of consecutive points
        # lat1, lon1 = data_KM.loc[i - 1, 'latitude'], data_KM.loc[i - 1, 'longitude']
        # lat2, lon2 = data_KM.loc[i, 'latitude'], data_KM.loc[i, 'longitude']
        # lat1, lon1 = data.loc[i - 1, 'latitude'], data.loc[i - 1, 'longitude']
        # lat2, lon2 = data.loc[i, 'latitude'], data.loc[i, 'longitude']
         # Get latitude and longitude of consecutive points
        # lat1 = data.iloc[i, data.columns.get_loc('latitude')]
        # lon1 = data.iloc[i, data.columns.get_loc('longitude')]
        # lat2 = data.iloc[i + 1, data.columns.get_loc('latitude')]
        # lon2 = data.iloc[i + 1, data.columns.get_loc('longitude')]
   
        # Get latitude and longitude of consecutive points
        lat1 = data['latitude'].iloc[i]
        lon1 = data['longitude'].iloc[i]
        lat2 = data['latitude'].iloc[i + 1]
        lon2 = data['longitude'].iloc[i + 1]
       
 
        # Calculate distance between consecutive points
        distance = haversine(lat1, lon1, lat2, lon2)
 
        # Add distance to total distance covered
        total_distance += distance
        # print("TotalDistancw----------->",total_distance)
 
    print("Total distance covered (in kilometers):{:.2f}".format(total_distance))
 
    ##############   Wh/Km
    Wh_km = abs(watt_h / total_distance)
    print("WH/KM:{:.2f}". format (watt_h / total_distance))
 
    # Assuming 'data' is your DataFrame with 'SOC_8' column
    initial_soc = data['SOC_8'].iloc[-1]  # Initial SOC percentage
    final_soc = data['SOC_8'].iloc[0]   # Final SOC percentage
 
    # Calculate total SOC consumed
    total_soc_consumed =  abs(final_soc - initial_soc)
 
    print ("Total SOC consumed:{:.2f}".format (total_soc_consumed),"%")
 
 
    # Check if the mode remains constant or changes
 
 
    #############################
    mode_values = data['Mode_Ack_408094978'].unique()
 
    if len(mode_values) == 1:
        # Mode remains constant throughout the log file
        mode = mode_values[0]
        if mode == 3:
            print("Mode is Custom mode.")
        elif mode == 2:
            print("Mode is Sports mode.")
        elif mode == 1:
            print("Mode is Eco mode.")
    else:
        # Mode changes throughout the log file
        mode_counts = data['Mode_Ack_408094978'].value_counts(normalize=True) * 100
        for mode, percentage in mode_counts.items():
            if mode == 3:
                print(f"Custom mode: {percentage:.2f}%")
            elif mode == 2:
                print(f"Sports mode: {percentage:.2f}%")
            elif mode == 1:
                print(f"Eco mode: {percentage:.2f}%")
    ###############################
 
 
  ##preak current##
    # Calculate power using PackCurr_6 and PackVol_6
    data_resampled['Power'] = data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
 
    # Find the peak power
    peak_power = data_resampled['Power'].min()
    print("Peak Power:", peak_power)
   
    # Calculate the average power
    average_power = abs(data_resampled['Power'].mean())
    print("Average Power:", average_power)
    altitude = None
    # Calculate altitude if 'Altitude' column exists in data_resampled
    if 'Altitude' in data_resampled.columns:
        altitude = data_resampled['Altitude'].max()  # Assuming you want to find the maximum altitude
        print("Maximum Altitude:", altitude)
    else:
        print("Altitude data not available.")
 
 
    # Calculate energy regenerated (in watt-hours)
    energy_regenerated = ((data_resampled[data_resampled['Power'] > 0]['Power']*data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
 
    # Calculate total energy consumed (in watt-hours)
    total_energy_consumed =  ((data_resampled[data_resampled['Power'] < 0]['Power']*data_resampled['Time_Diff']).sum()) / 3600  # Convert seconds to hours   ######################################################################
 
    print("total",total_energy_consumed)
    print("total energy regenerated",energy_regenerated)
 
   
   
   
    # Calculate regenerative effectiveness as a percentage
    if total_energy_consumed != 0:
     regenerative_effectiveness = abs(energy_regenerated / total_energy_consumed) * 100
     print("Regenerative Effectiveness (%):", regenerative_effectiveness)
    else:
     print("Total energy consumed is 0, cannot calculate regenerative effectiveness.")
 
    # Calculate idling time percentage (RPM was zero for more than 5 seconds)
    idling_time = ((data['MotorSpeed_340920578'] <= 0) | data['MotorSpeed_340920578'].isna()).sum()
    print("data length time",len(data))
    idling_percentage = (idling_time / len(data)) * 100
    print("Idling time percentage:", idling_percentage)
 
    # Calculate time spent in specific speed ranges
    speed_ranges = [(0, 10), (10, 20), (20, 30), (30, 40), (40, 50),(50, 60),(60,70),(70, 80),(80, 90)]
    speed_range_percentages = {}
 
    for range_ in speed_ranges:
        speed_range_time = ((data['MotorSpeed_340920578'] * 0.016 > range_[0]) & (data['MotorSpeed_340920578'] * 0.016 < range_[1])).sum()
        speed_range_percentage = (speed_range_time / len(data)) * 100
        speed_range_percentages[f"Time spent in {range_[0]}-{range_[1]} km/h"] = speed_range_percentage
        print(f"Time spent in {range_[0]}-{range_[1]} km/h: {speed_range_percentage:.2f}%")
 
 
##############################################################################################################################################################################
           
    # Calculate power using PackCurr_6 and PackVol_6
    data_resampled['Power'] = -data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
 
    # Find the peak power
    peak_power = data_resampled['Power'].max()
 
    # Get the maximum cell voltage
    max_cell_voltage = data_resampled['MaxCellVol_5'].max()
 
    # Find the index where the maximum voltage occurs
    max_index = data_resampled['MaxCellVol_5'].idxmax()
 
    # Retrieve the corresponding cell ID using the index
    max_cell_id = data_resampled['MaxVoltId_5'].loc[max_index]
 
    # Get the minimum cell voltage
    min_cell_voltage = data_resampled['MinCellVol_5'].min()
 
    # Find the index where the minimum voltage occurs
    min_index = data_resampled['MinCellVol_5'].idxmin()
 
    # Retrieve the corresponding cell ID using the index
    min_cell_id = data_resampled['MinVoltId_5'].loc[min_index]
 
    voltage_difference = max_cell_voltage - min_cell_voltage
 
    # Get the maximum temperature
    max_temp = data_resampled['MaxTemp_7'].max()
 
    # Find the index where the maximum temperature occurs
    max_temp_index = data_resampled['MaxTemp_7'].idxmax()
 
    # Retrieve the corresponding temperature ID using the index
    max_temp_id = data_resampled['MaxTempId_7'].loc[max_temp_index]
 
    # Get the minimum temperature
    min_temp = data_resampled['MinTemp_7'].min()
 
    # Find the index where the minimum temperature occurs
    min_temp_index = data_resampled['MinTemp_7'].idxmin()
 
    # Retrieve the corresponding temperature ID using the index
    min_temp_id = data_resampled['MinTempId_7'].loc[min_temp_index]
 
    # Calculate the difference in temperature
    temp_difference = max_temp - min_temp
    print()
 
    # Get the maximum temperature of FetTemp_8
    max_fet_temp = data_resampled['FetTemp_8'].max()
 
    # Get the maximum temperature of AfeTemp_12
    max_afe_temp = data_resampled['AfeTemp_12'].max()
 
    # Get the maximum temperature of PcbTemp_12
    max_pcb_temp = data_resampled['PcbTemp_12'].max()
 
    # Get the maximum temperature of MCU_Temperature_408094979
    max_mcu_temp = data_resampled['MCU_Temperature_408094979'].max()
 
    # Check for abnormal motor temperature at high RPMs
    max_motor_temp = data_resampled['Motor_Temperature_408094979'].max()
   
    # Find the battery voltage
    batteryVoltage = (data_resampled['BatteryVoltage_340920578'].max()) * 10
    print( "Battery Voltage", batteryVoltage )
 
    # Check for abnormal motor temperature at high RPMs for at least 15 seconds
    abnormal_motor_temp = (data_resampled['Motor_Temperature_408094979'] < 10) & (data_resampled['MotorSpeed_340920578'] > 3500)
    abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
 
    # Check if abnormal condition persists for at least 15 seconds
    abnormal_motor_temp_detected = (abnormal_motor_temp_mask >= 120).any()
 
    #For battery Analysis
    cycleCount= data_resampled['CycleCount_7'].max()
 
   
 
    ByteBeamId= data['id'].iloc[0]
   
 
 
    # # Filter rows where SOC is 100%
    # data_100_soc = data_resampled[data_resampled['SOC_8'] == 100]
 
    # # Find the lowest cell temperature at 100% SOC
    # lowest_temp_100_soc = data_100_soc['Temp1_10'].min()
    # print("Lowest Cell Temperature at 100% SOC:", lowest_temp_100_soc)
 
    # # Filter rows where SOC is 100%
    # data_100_soc = data_resampled[data_resampled['SOC_8'] == 100]
 
    # # Find the highest cell temperature at 100% SOC
    # highest_temp_100_soc = data_100_soc['Temp1_10'].max()
    # print("highest Cell Temperature at 100% SOC:", highest_temp_100_soc)
 
    # temp_difference = highest_temp_100_soc - lowest_temp_100_soc
    # print("Difference between Highest and Lowest Cell Temperature at 100% SOC:", temp_difference)
 
    # Find the maximum BMS temperature in Celsius
    #max_bms_temp = data_resampled['FetTemp_8'].max()
    #print("Maximum BMS Temperature in Celsius:", max_bms_temp)
 
    total_energy_kwh = actual_ah * batteryVoltage / 1000
    print("Total energy charged in kWh: {:.2f}".format(total_energy_kwh))
 
    total_energy_kw = total_energy_kwh / total_duration.seconds / 3600
    print("Electricity consumption units in kW", (total_energy_kw))
 
    # Add these variables and logic to ppt_data
    ppt_data = {
        "Date and Time": str(start_time_seconds) + " to " + str(end_time_seconds),
        "Byte Beam ID": ByteBeamId,
        "Total time taken for the ride": total_duration,
        "Actual Ampere-hours (Ah)": actual_ah,
        "Actual Watt-hours (Wh)- Calculated_UsingFormala": watt_h,
        "Starting SoC (Ah)": starting_soc_Ah,
        "Ending SoC (Ah)": ending_soc_Ah,
        "Starting SoC (%)": starting_soc_percentage,
        "Ending SoC (%)": ending_soc_percentage,
        "Total distance covered (km)": total_distance,
        "Total energy consumption(WH/KM)": watt_h / total_distance,
        "Total SOC consumed(%)":starting_soc_percentage- ending_soc_percentage,
        "Mode": "",
        "Peak Power(kW)": peak_power,
        "Average Power(kW)": average_power,
        "Total Energy Regenerated(kWh)": energy_regenerated,
        "Regenerative Effectiveness(%)": regenerative_effectiveness,
        "Highest Cell Voltage(V)": max_cell_voltage,
        "Lowest Cell Voltage(V)": min_cell_voltage,
        "Difference in Cell Voltage(V)": voltage_difference,
        "Minimum Temperature(C)": min_temp,
        "Maximum Temperature(C)": max_temp,
        "Difference in Temperature(C)": max_temp- min_temp,
 
        "Maximum Fet Temperature-BMS(C)": max_fet_temp,
        "Maximum Afe Temperature-BMS(C)": max_afe_temp,
        "Maximum PCB Temperature-BMS(C)": max_pcb_temp,
        "Maximum MCU Temperature(C)": max_mcu_temp,
        "Maximum Motor Temperature(C)": max_motor_temp,
        "Abnormal Motor Temperature Detected(C)": abnormal_motor_temp_detected,
        "highest cell temp(C)": max_cell_temp,
        "lowest cell temp(C)": min_cell_temp,
        "Difference between Highest and Lowest Cell Temperature at 100% SOC(C)": CellTempDiff,
        "Battery Voltage(V)": batteryVoltage,
        "Total energy charged(kWh)- Calculated_BatteryData": total_energy_kwh,
        "Electricity consumption units(kW)": total_energy_kw,
        "Cycle Count of battery": cycleCount
        }
 
   ######################################
    ######################################
 
    mode_values = data_resampled['Mode_Ack_408094978'].unique()
    if len(mode_values) == 1:
        mode = mode_values[0]
        if mode == 3:
            ppt_data["Mode"] = "Custom mode"
        elif mode == 2:
            ppt_data["Mode"] = "Sports mode"
        elif mode == 1:
            ppt_data["Mode"] = "Eco mode"
    else:
        # Mode changes throughout the log file
        mode_counts = data_resampled['Mode_Ack_408094978'].value_counts(normalize=True) * 100
        mode_strings = []  # Initialize list to store mode strings
        for mode, percentage in mode_counts.items():
            if mode == 3:
                mode_strings.append(f"Custom mode\n{percentage:.2f}%")
            elif mode == 2:
                mode_strings.append(f"Sports mode\n{percentage:.2f}%")
            elif mode == 1:
                mode_strings.append(f"Eco mode\n{percentage:.2f}%")
        ppt_data["Mode"] = "\n".join(mode_strings)
 
     # Add calculated parameters to ppt_data
    ppt_data["Idling time percentage"] = idling_percentage
    ppt_data.update(speed_range_percentages)
################################################################################################# recent data
 
    # Calculate power using PackCurr_6 and PackVol_6
    data_resampled['Power'] = -data_resampled['PackCurr_6'] * data_resampled['PackVol_6']
 
    # Find the peak power
    peak_power = data_resampled['Power'].max()
    print("Peak Power:", peak_power)
 
    # Get the maximum cell voltage
    max_cell_voltage = data_resampled['MaxCellVol_5'].max()
 
    # Find the index where the maximum voltage occurs
    max_index = data_resampled['MaxCellVol_5'].idxmax()
 
    # Retrieve the corresponding cell ID using the index
    max_cell_id = data_resampled['MaxVoltId_5'].loc[max_index]
 
        # Get the minimum cell voltage
    min_cell_voltage = data_resampled['MinCellVol_5'].min()
 
    # Find the index where the minimum voltage occurs
    min_index = data_resampled['MinCellVol_5'].idxmin()
 
    # Retrieve the corresponding cell ID using the index
    min_cell_id = data_resampled['MinVoltId_5'].loc[min_index]
 
    voltage_difference = max_cell_voltage - min_cell_voltage
 
 
    print("Lowest Cell Voltage:", min_cell_voltage, "V, Cell ID:", min_cell_id)
    print("Highest Cell Voltage:", max_cell_voltage, "V, Cell ID:", max_cell_id)
    print("Difference in Cell Voltage:", voltage_difference, "V")
 
    # Get the maximum temperature
    max_temp = data_resampled['MaxTemp_7'].max()
 
    # Find the index where the maximum temperature occurs
    max_temp_index = data_resampled['MaxTemp_7'].idxmax()
 
    # Retrieve the corresponding temperature ID using the index
    max_temp_id = data_resampled['MaxTempId_7'].loc[max_temp_index]
 
 
    # Get the minimum temperature
    min_temp = data_resampled['MinTemp_7'].min()
 
    # Find the index where the minimum temperature occurs
    min_temp_index = data_resampled['MinTemp_7'].idxmin()
 
    # Retrieve the corresponding temperature ID using the index
    min_temp_id = data_resampled['MinTempId_7'].loc[min_temp_index]
    # Calculate the difference in temperature
    temp_difference = max_temp - min_temp
    print("temp_difference------------------------>",temp_difference)
    print("Total distance------------------------->",total_distance)
 
    # Print the information
    print("Maximum Temperature:", max_temp, "C, Temperature ID:", max_temp_id)
    print("Minimum Temperature:", min_temp, "C, Temperature ID:", min_temp_id)
    print("Difference in Temperature:", temp_difference, "C")
 
    # Get the maximum temperature of FetTemp_8
    max_fet_temp = data_resampled['FetTemp_8'].max()
    print("Maximum Fet Temperature:", max_fet_temp, "C")
 
    # Get the maximum temperature of AfeTemp_12
    max_afe_temp = data_resampled['AfeTemp_12'].max()
    print("Maximum Afe Temperature:", max_afe_temp, "C")
 
    # Get the maximum temperature of PcbTemp_12
    max_pcb_temp = data_resampled['PcbTemp_12'].max()
    print("Maximum PCB Temperature:", max_pcb_temp, "C")
 
    # Get the maximum temperature of MCU_Temperature_408094979
    max_mcu_temp = data_resampled['MCU_Temperature_408094979'].max()
    print("Maximum MCU Temperature:", max_mcu_temp, "C")
 
    # Check for abnormal motor temperature at high RPMs
    max_motor_temp = data_resampled['Motor_Temperature_408094979'].max()
 
 
    print("Maximum Motor Temperature:", max_motor_temp, "C")
 
 
    # Check for abnormal motor temperature at high RPMs for at least 15 seconds
    abnormal_motor_temp = (data_resampled['Motor_Temperature_408094979'] < 10) & (data_resampled['MotorSpeed_340920578'] > 3500)
   
    # Convert to a binary mask indicating consecutive occurrences
    abnormal_motor_temp_mask = abnormal_motor_temp.astype(int).groupby(abnormal_motor_temp.ne(abnormal_motor_temp.shift()).cumsum()).cumsum()
 
    #Check if abnormal condition persists for at least 15 seconds
    #if (abnormal_motor_temp_mask >= 15).any():
    #print("Abnormal motor temperature detected: NTC has issues - very low temperature at high RPMs")
 
 
    return total_duration, total_distance, Wh_km, total_soc_consumed,ppt_data
 
 
 
# folder_path = r"C:\Users\kamalesh.kb\CodeForAutomation\MAIN_FOLDER\MAR_21"
 
 
####################
   
# def capture_analysis_output(log_file, km_file, folder_path):
def capture_analysis_output(log_file,folder_path):
    try:
        # Capture print statements
        analysis_output = io.StringIO()
        output_file = "analysis_results.docx"
 
        with redirect_stdout(analysis_output):
            #total_duration, total_distance, Wh_km, total_soc_consumed, ppt_data = analysis_Energy(log_file, km_file)
            analysis_output = analysis_output.getvalue()
 
        # Extract folder name from folder_path
        folder_name = os.path.basename(folder_path)
 
        # Create a new PowerPoint presentation
        prs = Presentation()
 
        # Add title slide with 'Selawik' style
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        title.text = f"Analysis Results from Folder - {folder_name}"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)  # Corrected to Pt
        title.text_frame.paragraphs[0].font.name = 'Selawik'
 
        rows = len(ppt_data) + 1
        cols = 2
        table_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(table_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = "Analysis Results:"
 
        # Centering the title horizontally
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.5)  # Adjust as needed
 
        # Setting the font size of the title
        title_shape.text_frame.paragraphs
 
        # Define maximum number of rows per slide
        max_rows_per_slide = 13
        # Add some space between title and table
        title_shape.top = Inches(0.5)
 
        table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
        table.columns[0].width = Inches(4)
        table.columns[1].width = Inches(4)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
 
        # Initialize row index
        row_index = 1
 
        # Iterate over data and populate the table
        for key, value in ppt_data.items():
            # Check if current slide has reached maximum rows
            if row_index > max_rows_per_slide:
                # Add a new slide
                slide = prs.slides.add_slide(table_slide_layout)
                shapes = slide.shapes
                title_shape = shapes.title
                title_shape.text = "Analysis Results:"
 
                # Add a new table to the new slide
                table = shapes.add_table(max_rows_per_slide + 1, cols, Inches(1), Inches(1.5), Inches(8), Inches(5)).table
                table.columns[0].width = Inches(4)
                table.columns[1].width = Inches(4)
                table.cell(0, 0).text = "Metric"
                table.cell(0, 1).text = "Value"
 
                # Reset row index for the new slide
                row_index = 1
 
            # Populate the table
            table.cell(row_index, 0).text = key
            table.cell(row_index, 1).text = str(value)
 
            # Increment row index
            row_index += 1
 
        # Add image slide with title and properly scaled image
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
 
        # Remove the unwanted title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                slide.shapes._spTree.remove(shape._element)
 
        # Add the title
        title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.5), prs.slide_width - Inches(2), Inches(1))
        title_shape.text = "Graph Analysis"
 
        # Add the image and adjust its position and size
        graph_width = prs.slide_width - Inches(1)
        graph_height = prs.slide_height - Inches(2)
        left = (prs.slide_width - graph_width) / 2
        top = (prs.slide_height - graph_height) / 2 + Inches(1)
        pic = slide.shapes.add_picture('graph.png', left, top, width=graph_width, height=graph_height)
 
        # Save the presentation
        output_file_name = f"{folder_path}/analysis_{folder_name}.pptx"
        prs.save(output_file_name)
        print("PPT generated!")
 
        # Create a new Excel workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Analysis Results"
 
        # Populate Excel sheet with analysis data
        for i, (key, value) in enumerate(ppt_data.items(), start=1):
            ws.cell(row=i, column=1, value=key)
            ws.cell(row=i, column=2, value=value)
 
        # Save the Excel workbook
        excel_output_file = f"{folder_path}/analysis_{folder_name}.xlsx"
        wb.save(excel_output_file)
        print("Excel generated!")
 
    except Exception as e:
        print("Error:", e)
 
 
 
 
 
####################
 
#folder_path = "/home/sanjith/Documents/Graphs _ creta/15-51_16-00"
 
# Get the list of files in the folder
#files = os.listdir(folder_path)
 
# Initialize variables to store file paths
log_file = None
# km_file = None
 
 
 
main_folder_path = r"D:\March_BB4\March_BB4\Mar-30"

 
def mergeExcel(main_folder_path):
    def prepare_sheet_in_memory(file_path):
        workbook = load_workbook(filename=file_path)
        sheet = workbook.active
        file_name = os.path.basename(file_path)
        sheet.insert_rows(1)
        sheet['A1'] = 'file name'
        sheet['B1'] = file_name
        return sheet, file_name
 
    def sheet_to_dict(sheet):
        data_dict = {}
        for row in sheet.iter_rows(min_row=2, values_only=True):
            key, *values = row
            data_dict[key] = values
        return data_dict
 
    def merge_dicts(dict1, dict2):
        for key, values in dict2.items():
            if key in dict1:
                # Assuming that values is a list of values
                dict1[key].extend(values)
            else:
                dict1[key] = values
        return dict1
 
    def process_directory(directory):
        merged_data = {}
        file_names = []
        for root, dirs, files in os.walk(directory):
            for name in dirs:
                if name.startswith("B"):
                    subdir_path = os.path.join(root, name)
                    for file_name in os.listdir(subdir_path):
                        if file_name.endswith(".xlsx"):
                            file_path = os.path.join(subdir_path, file_name)
                            sheet, extracted_file_name = prepare_sheet_in_memory(file_path)
                            data_dict = sheet_to_dict(sheet)
                            merged_data = merge_dicts(merged_data, data_dict)
                            if extracted_file_name not in file_names:
                                file_names.append(extracted_file_name)
        return merged_data, file_names
 
    def main(main_folder_path):
        directory = main_folder_path
        merged_data, file_names = process_directory(directory)
 
        merged_workbook = Workbook()
        merged_sheet = merged_workbook.active
 
        # Use the extracted file names for headers
        headers = ['File name'] + file_names
        merged_sheet.append(headers)
 
        if merged_data:
            for key, values in merged_data.items():
                merged_sheet.append([key] + values)
        else:
            print("No data found in merged_data")
 
        merged_file_path = os.path.join(directory, 'merged_analysis.xlsx')
        merged_workbook.save(filename=merged_file_path)
        print("Merged Excel file is ready")
 
    if __name__ == '__main__':
        main(main_folder_path)
 

 
 
for subfolder in os.listdir(main_folder_path):
    if subfolder.startswith("Battery"):  # Check if the subfolder starts with "Battery"
        subfolder_path = os.path.join(main_folder_path, subfolder)
        print(subfolder)
        if os.path.isdir(subfolder_path):
            log_file = None
            # km_file = None
            # Find 'log' and 'km' files
            log_found = False
            # km_found = False
            for file in os.listdir(subfolder_path):
                if file.startswith('log') and file.endswith('.csv'):
                    log_file = os.path.join(subfolder_path, file)
                    log_found = True
                # elif file.startswith('km') and file.endswith('.csv'):
                #     km_file = os.path.join(subfolder_path, file)
                #     km_found = True
                # if log_found and km_found:
                if log_found:
                    break
 
            # if log_found and km_found:
            if log_found:
                data = pd.read_csv(log_file)
                # print("Read log files
                # data_KM = pd.read_csv(km_file)
                # print("Read log file------------------------>")
                total_duration = 0
                total_distance = 0
                Wh_km = 0
                SOC_consumed = 0
                mode_values = 0
 
                # Plot graphs
                plot_ghps(data)
                # total_duration, total_distance, Wh_km, SOC_consumed, ppt_data = analysis_Energy(log_file, km_file)
                # capture_analysis_output(log_file, km_file, subfolder_path)
                total_duration, total_distance, Wh_km, SOC_consumed, ppt_data = analysis_Energy(log_file )
                capture_analysis_output(log_file, subfolder_path,)
                
 
 
 
            else:
                print("Log file or KM file not found in subfolder:", subfolder)
 
 
mergeExcel(main_folder_path) 
                
 